import logging
import os
import random
import re
import sys
import time
from datetime import datetime

import matplotlib.pyplot as plt
import pandas as pd
import requests
import seaborn as sns
from PyQt5.QtCore import Qt, QTimer, QThread, pyqtSignal, QRectF, QSize, QDateTime
from PyQt5.QtGui import QFont, QPainterPath, QRegion, QPixmap, QIcon
from PyQt5.QtWidgets import (QApplication, QMainWindow, QWidget, QVBoxLayout, QHBoxLayout, QLabel, QPushButton,
                             QTextEdit, QLineEdit, QComboBox, QFrame, QTextBrowser, QStackedWidget,
                             QRadioButton, QMessageBox, QProgressBar, QFileDialog, QSizePolicy)
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# 配置中文字体
plt.rcParams['font.family'] = 'SimHei'
plt.rcParams['axes.unicode_minus'] = False

# DeepSeek API 配置
DEEPSEEK_API_KEY = "sk-986dc7e18c9442aaac66207428d28600"
SAFE_CHARS = set("!@$%^&()_-+=[]{}',.`~")
MAX_RETRIES = 3
TIMEOUT = 80

def call_deepseek_api(prompt):
    headers = {"Authorization": f"Bearer {DEEPSEEK_API_KEY}", "Content-Type": "application/json"}
    data = {"messages": [{"role": "user", "content": prompt}], "model": "deepseek-chat", "temperature": 0.5}
    for attempt in range(MAX_RETRIES):
        try:
            response = requests.post("https://api.deepseek.com/v1/chat/completions", headers=headers, json=data,
                                     timeout=TIMEOUT)
            response.raise_for_status()
            result = response.json()
            content = result['choices'][0]['message']['content']
            return content
        except requests.Timeout:
            if attempt == MAX_RETRIES - 1:
                raise Exception("API请求超时，请检查网络连接")
            time.sleep(2 ** attempt)
        except Exception as e:
            raise Exception(f"API调用失败: {str(e)}")

# 日志配置
logging.basicConfig(
    filename=f'app_{datetime.now().strftime("%Y%m%d")}.log',
    level=logging.DEBUG,
    format='%(asctime)s - %(levelname)s - %(message)s',
    encoding='utf-8'
)

# 模板映射
TEMPLATE_MAP = {
    "教学课件类": "template_edu.pptx",
    "政务报告类": "template_gov.pptx",
    "学术答辩类": "template_academic.pptx",
    "简约类": "template_simple.pptx",
    "东软类": "template_neusoft.pptx"
}

class ChartWorker(QThread):
    status_update = pyqtSignal(str)
    finished = pyqtSignal(str)
    error = pyqtSignal(str)

    def __init__(self, file_path, chart_type, output_path):
        super().__init__()
        self.file_path = file_path
        self.chart_type = chart_type
        self.output_path = output_path
        self.cancelled = False

    def run(self):
        try:
            self.status_update.emit("正在读取 Excel 文件...")
            df = pd.read_excel(self.file_path)
            if df.empty:
                raise ValueError("Excel 文件为空")
            self.status_update.emit("正在生成图表...")
            plt.figure(figsize=(10, 6))
            if self.chart_type == "饼状图":
                df.iloc[:, 1].value_counts().plot(kind='pie', autopct='%1.1f%%')
                plt.title(f"{df.columns[1]} 分布")
            elif self.chart_type == "柱状图":
                df.plot(kind='bar', x=df.columns[0], y=df.columns[1])
                plt.title(f"{df.columns[1]} 柱状图")
            elif self.chart_type == "折线图":
                df.plot(kind='line', x=df.columns[0], y=df.columns[1])
                plt.title(f"{df.columns[1]} 趋势")
            elif self.chart_type == "散点图":
                df.plot(kind='scatter', x=df.columns[0], y=df.columns[1])
                plt.title(f"{df.columns[0]} vs {df.columns[1]} 散点图")
            elif self.chart_type == "热力图":
                corr_matrix = df.select_dtypes(include=['number']).corr()
                sns.heatmap(corr_matrix, annot=True, cmap='coolwarm')
                plt.title("热力图 - 相关性矩阵")
            elif self.chart_type == "直方图":
                df[df.columns[1]].plot(kind='hist', bins=10, edgecolor='black')
                plt.title(f"{df.columns[1]} 分布直方图")
            elif self.chart_type == "气泡图":
                plt.scatter(df.iloc[:, 0], df.iloc[:, 1], s=df.iloc[:, 2] * 50, alpha=0.6)
                plt.title(f"气泡图：{df.columns[0]} vs {df.columns[1]}")
            plt.tight_layout()
            if self.cancelled:
                self.finished.emit("图表生成已取消")
                return
            self.status_update.emit("正在保存图表...")
            file_ext = os.path.splitext(self.output_path)[1][1:].lower()
            plt.savefig(self.output_path, dpi=300, format=file_ext)
            plt.close()
            self.finished.emit(f"图表已保存至:\n{self.output_path}")
        except Exception as e:
            self.error.emit(f"生成图表失败: {str(e)}")

    def cancel(self):
        self.cancelled = True

class PPTGenerator:
    def __init__(self, update_status=None):
        self.prs = None
        self.layout_map = {'cover': 0, 'section': 1, 'content': [2, 3], 'chart': 3, 'quote': 4, 'end': 4}
        self.available_content_layouts = [2, 3]
        self.cancelled = False
        self.update_status = update_status
        self.start_time = None
        self.current_stage = "就绪"
        self.selected_template = "template_simple.pptx"

    def set_template(self, template_name):
        self.selected_template = TEMPLATE_MAP.get(template_name, "template_simple.pptx")
        if not os.path.exists(self.selected_template):
            logging.error(f"模板文件不存在: {self.selected_template}")
            raise FileNotFoundError(f"模板文件 {self.selected_template} 不存在")
        logging.info(f"用户选择模板: {template_name} -> {self.selected_template}")

    def safe_filename(self, filename):
        return ''.join(c for c in filename if c.isalnum() or c in SAFE_CHARS)

    def clean_markdown(self, text):
        text = re.sub(r'[#\*\>\-]+\s*', '', text)
        return text.strip()

    def call_deepseek_api(self, prompt):
        headers = {"Authorization": f"Bearer {DEEPSEEK_API_KEY}", "Content-Type": "application/json"}
        data = {"messages": [{"role": "user", "content": prompt}], "model": "deepseek-chat", "temperature": 0.5}
        for attempt in range(MAX_RETRIES):
            try:
                response = requests.post("https://api.deepseek.com/v1/chat/completions", headers=headers, json=data,
                                         timeout=TIMEOUT)
                response.raise_for_status()
                result = response.json()
                content = result['choices'][0]['message']['content']
                logging.debug(f"API返回内容长度: {len(content)} 字符, 内容: {content[:100]}...")
                return content[:1500] if len(content) > 1500 else content
            except requests.Timeout:
                logging.warning(f"API请求超时，第 {attempt + 1}/{MAX_RETRIES} 次尝试")
                if attempt == MAX_RETRIES - 1:
                    raise Exception("API请求超时，请检查网络连接")
                time.sleep(2 ** attempt)
            except Exception as e:
                logging.error(f"API调用失败: {str(e)}")
                raise Exception(f"API调用失败: {str(e)}")

    def generate_outline(self, topic):
        self.current_stage = "正在生成大纲"
        if self.update_status:
            self.update_status(self.current_stage)
        prompt = (
            f"生成关于【{topic}】的PPT大纲，格式为Markdown，包含4-6个章节，每个章节3-5个要点，"
            f"要求内容详细且丰富，提供具体的描述或示例，插入数据标记如{{data:market_size}}"
        )
        outline = self.call_deepseek_api(prompt)
        logging.info(f"生成的大纲长度: {len(outline)} 字符, 大纲: {outline}")
        return outline

    def expand_content(self, section_title, content):
        self.current_stage = f"正在扩展章节内容: {section_title}"
        if self.update_status:
            self.update_status(self.current_stage)
        content = re.sub(re.escape(section_title), '', content, flags=re.IGNORECASE).strip()
        prompt = (
            f"请为PPT章节【{section_title}】生成详细内容，基于以下现有内容：\n"
            f"```\n{content}\n```\n"
            f"严格按照以下Markdown格式返回：\n"
            f"```\n"
            f"## 正文\n"
            f"基于现有内容优化或补充，生成100-200字的详细正文，要求内容丰富，包含具体细节、示例或数据（如市场规模、趋势分析等），按段落分隔。\n"
            f"## 备注\n"
            f"生成与正文相关的拓展内容，必须包含以下两部分：\n"
            f"- 知识点：列出2-3个与正文相关的核心知识点（每点一句）。\n"
            f"- 提问：列出2-3个适合课堂提问的简单问题（每题一句）。\n"
            f"```\n"
            f"示例：\n"
            f"```\n"
            f"## 正文\n"
            f"机器学习是AI的核心技术，通过数据训练模型预测结果。例如，垃圾邮件过滤器利用分类算法，准确率达95%。\n"
            f"## 备注\n"
            f"- 知识点：\n"
            f"  - 机器学习基于统计学和概率论。\n"
            f"  - 深度学习是神经网络的一种。\n"
            f"- 提问：\n"
            f"  - 机器学习如何通过数据学习？\n"
            f"  - 深度学习有哪些应用？\n"
            f"```\n"
            f"确保备注部分不可为空，且与正文内容相关！"
        )
        try:
            expanded = self.call_deepseek_api(prompt)
            logging.info(f"扩展内容总长度: {len(expanded)} 字符, 原始返回: {expanded[:150]}...")
            main_content_match = re.search(r'## 正文\s*([\s\S]*?)(?=## 备注|$)', expanded, re.DOTALL)
            notes_content_match = re.search(r'## 备注\s*([\s\S]*)', expanded, re.DOTALL)
            main_content = main_content_match.group(1).strip() if main_content_match else content
            notes_content = notes_content_match.group(1).strip() if notes_content_match else None

            if not notes_content or "知识点" not in notes_content or "提问" not in notes_content:
                logging.warning(f"API返回的备注格式不完整: {notes_content}")
                notes_content = (
                    "知识点：\n"
                    "- 此处为默认知识点1。\n"
                    "- 此处为默认知识点2。\n"
                    "提问：\n"
                    "- 此处为默认问题1？\n"
                    "- 此处为默认问题2？"
                )

            main_content = re.sub(r'{data:\w+}', 'X', main_content)
            notes_content = self.clean_markdown(notes_content)
            logging.debug(f"正文长度: {len(main_content)} 字符, 备注长度: {len(notes_content)} 字符")
            return main_content, notes_content
        except Exception as e:
            logging.error(f"扩展内容失败: {str(e)}")
            default_main = re.sub(r'{data:\w+}', 'X', content)
            default_notes = (
                "知识点：\n"
                "- 此处为默认知识点1。\n"
                "- 此处为默认知识点2。\n"
                "提问：\n"
                "- 此处为默认问题1？\n"
                "- 此处为默认问题2？"
            )
            return default_main, default_notes

    def parse_input(self, input_text):
        self.current_stage = "正在解析输入"
        if self.update_status:
            self.update_status(self.current_stage)
        if not input_text.strip():
            raise ValueError("输入内容不能为空")
        if re.search(r'^#+\s|\* |- |\d+\.', input_text, re.MULTILINE):
            return self.parse_markdown(input_text)
        else:
            return self.parse_markdown(self.generate_outline(input_text))

    def parse_markdown(self, md_text):
        data = {"title": "", "sections": []}
        current_section = None
        current_paragraph = []

        for line in md_text.split('\n'):
            line = line.strip()
            if not line:
                if current_paragraph and current_section:
                    current_section["content"].append(self.clean_markdown('\n'.join(current_paragraph)))
                    current_paragraph = []
                continue

            if line.startswith('# '):
                data["title"] = self.clean_markdown(line)
                logging.debug(f"解析封面标题: {data['title']}")
            elif line.startswith('## '):
                if current_section:
                    if current_paragraph:
                        current_section["content"].append(self.clean_markdown('\n'.join(current_paragraph)))
                        current_paragraph = []
                    data["sections"].append(current_section)
                cleaned_title = self.clean_markdown(line)
                current_section = {"title": cleaned_title, "content": []}
                logging.debug(f"解析二级标题: {cleaned_title}")
            elif re.match(r'^[\-\*] ', line) and current_section:
                if current_paragraph:
                    current_section["content"].append(self.clean_markdown('\n'.join(current_paragraph)))
                    current_paragraph = []
                cleaned_content = self.clean_markdown(line)
                if cleaned_content and cleaned_content != current_section["title"]:
                    current_section["content"].append(cleaned_content)
            elif re.match(r'^>\s+', line) and current_section:
                if current_paragraph:
                    current_section["content"].append(self.clean_markdown('\n'.join(current_paragraph)))
                    current_paragraph = []
                current_section["content"].append(self.clean_markdown(line))
            elif current_section:
                current_paragraph.append(line)

        if current_section:
            if current_paragraph:
                current_section["content"].append(self.clean_markdown('\n'.join(current_paragraph)))
            data["sections"].append(current_section)

        if not data["title"] or not data["sections"]:
            raise ValueError("输入必须包含标题和至少一个章节")
        return data

    def build_ppt(self, data, output_path):
        if self.cancelled:
            logging.info("生成被用户取消")
            return None
        self.start_time = time.time()
        output_path = os.path.abspath(output_path)
        logging.info(f"构建PPT时传入的路径: {output_path}")
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        self._init_presentation()
        self._create_cover_slide(data)
        self._create_toc_slide(data)
        for idx, section in enumerate(data["sections"]):
            if self.cancelled:
                logging.info("生成被用户取消")
                return None
            layout_type = self._detect_layout_type(section)
            self._add_section_slides(layout_type, section, idx)
        self._add_end_slide()
        self._save_ppt(output_path)
        return output_path

    def _init_presentation(self):
        self.current_stage = "正在初始化PPT"
        if self.update_status:
            self.update_status(self.current_stage)
        try:
            self.prs = Presentation(self.selected_template)
            while self.prs.slides:
                self.prs.part.drop_rel(self.prs.slides._sldIdLst[-1].rId)
                del self.prs.slides._sldIdLst[-1]
        except Exception as e:
            logging.warning(f"加载模板 {self.selected_template} 失败: {str(e)}，使用默认设置")
            self.prs = Presentation()
            self.available_content_layouts = [2, 3]
            self.prs.slide_width = Inches(13.33)
            self.prs.slide_height = Inches(7.5)

    def _create_cover_slide(self, data):
        self.current_stage = "生成封面页"
        if self.update_status:
            self.update_status(self.current_stage)
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[self.layout_map['cover']])
        title = slide.shapes.title or slide.shapes.add_textbox(Inches(1), Inches(1), Inches(11), Inches(2))
        title.text = data["title"]
        p = title.text_frame.paragraphs[0]
        p.font.size = Pt(44)
        p.font.name = '微软雅黑'
        p.alignment = PP_ALIGN.CENTER
        subtitle = None
        if len(slide.placeholders) > 1:
            for ph in slide.placeholders:
                if ph.placeholder_format.idx != 0:
                    subtitle = ph
                    break
        if not subtitle:
            logging.info(f"封面页缺少副标题占位符，自动添加文本框")
            subtitle = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11), Inches(1))
        subtitle.text = "吴卓宪同学制作的AI智能体制作"
        p = subtitle.text_frame.paragraphs[0]
        p.font.size = Pt(24)
        p.font.name = '微软雅黑'
        p.font.color.rgb = RGBColor(0x42, 0x85, 0xF4)
        p.alignment = PP_ALIGN.CENTER
        logging.debug(f"封面页标题: {data['title']}")

    def _create_toc_slide(self, data):
        self.current_stage = "生成目录页"
        if self.update_status:
            self.update_status(self.current_stage)
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[2])
        title = slide.shapes.title or slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))
        title.text = "目录"
        p = title.text_frame.paragraphs[0]
        p.font.name = '微软雅黑'
        p.font.size = Pt(36)
        p.alignment = PP_ALIGN.CENTER

        content_box = None
        if len(slide.placeholders) > 1:
            for ph in slide.placeholders:
                if ph.placeholder_format.idx != 0:
                    content_box = ph
                    break
        if not content_box:
            logging.info(f"目录页缺少正文占位符，自动添加文本框")
            content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(4))

        tf = content_box.text_frame
        tf.clear()
        tf.word_wrap = True
        for idx, section in enumerate(data["sections"], 1):
            p = tf.add_paragraph()
            p.text = f"{idx}. {section['title']}"
            p.font.size = Pt(20)
            p.font.name = '微软雅黑'
            p.space_after = Pt(10)
            p.level = 0

    def _detect_layout_type(self, section):
        content = "\n".join(section.get("content", []))
        if "{chart:" in content:
            return 'chart'
        if re.search(r'^>\s+', content, re.MULTILINE):
            return 'quote'
        if len(section.get("content", [])) < 2:
            return 'section'
        return 'content'

    def _create_slide(self, layout_type):
        try:
            if layout_type == 'content':
                layout_idx = random.choice(self.available_content_layouts)
            else:
                layout_idx = self.layout_map[layout_type] if isinstance(self.layout_map[layout_type], int) else \
                self.layout_map[layout_type][0]
            slide = self.prs.slides.add_slide(self.prs.slide_layouts[layout_idx])
            logging.debug(f"创建幻灯片 {layout_type}，使用版式索引 {layout_idx}，占位符数: {len(slide.placeholders)}")
            return slide
        except (IndexError, KeyError) as e:
            logging.warning(f"布局 {layout_type} 不可用: {str(e)}，使用默认布局")
            return self.prs.slides.add_slide(self.prs.slide_layouts[1])

    def _add_section_slides(self, layout_type, section, idx):
        self.current_stage = f"生成章节: {idx + 1}. {section['title']}"
        if self.update_status:
            self.update_status(self.current_stage)
        content = "\n".join(section.get("content", [])) or f"这是关于 {section['title']} 的默认内容。"
        main_content, notes_content = self.expand_content(section['title'], content)
        logging.info(
            f"章节 {idx + 1}. {section['title']} 正文长度: {len(main_content)} 字符, 备注长度: {len(notes_content)} 字符")

        if layout_type in ['chart', 'quote']:
            slide = self._create_slide(layout_type)
            self._add_slide_content(slide, section, idx, layout_type, main_content, notes_content=notes_content)
        else:
            content_parts = self._split_content(main_content)
            for i, part in enumerate(content_parts, 1):
                if self.cancelled:
                    return
                slide = self._create_slide(layout_type)
                title_suffix = f" ({i}/{len(content_parts)})" if len(content_parts) > 1 else ""
                self._add_slide_content(slide, section, idx, layout_type, part, title_suffix,
                                        notes_content if i == 1 else "（续前页备注）")

    def _split_content(self, content):
        if len(content) <= 800:
            return [content]
        parts = []
        current_part = ""
        for line in content.split('\n'):
            if len(current_part) + len(line) + 1 > 800 and current_part:
                parts.append(current_part.strip())
                current_part = line
            else:
                current_part += f"\n{line}" if current_part else line
        if current_part:
            parts.append(current_part.strip())
        return parts

    def _add_slide_content(self, slide, section, idx, layout_type, content, title_suffix="", notes_content=None):
        title = slide.shapes.title or slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))
        title.text = f"{idx + 1}. {section['title']}{title_suffix}"
        p = title.text_frame.paragraphs[0]
        p.font.name = '微软雅黑'
        p.font.size = {'cover': Pt(44), 'section': Pt(36), 'content': Pt(32), 'chart': Pt(28), 'quote': Pt(24)}.get(layout_type, Pt(32))
        p.alignment = PP_ALIGN.CENTER if layout_type in ['section', 'quote'] else PP_ALIGN.LEFT

        if layout_type == 'chart':
            self._add_chart(slide, content)
        elif layout_type == 'quote':
            self._add_quote(slide, content)
        else:
            self._add_normal_content(slide, content, layout_type)

        if notes_content:
            try:
                notes_slide = slide.notes_slide
                notes_tf = notes_slide.notes_text_frame
                notes_tf.clear()
                notes_tf.text = notes_content
                for p in notes_tf.paragraphs:
                    p.font.size = Pt(12)
                    p.font.name = '微软雅黑'
                logging.debug(f"幻灯片 {slide.slide_id} 添加备注成功: {notes_content[:50]}...")
            except Exception as e:
                logging.error(f"添加备注失败: {str(e)}")
                notes_tf.text = "备注添加失败，请检查模板或内容"

    def _add_normal_content(self, slide, content, layout_type):
        content_box = None
        if len(slide.placeholders) > 1:
            for ph in slide.placeholders:
                if ph.placeholder_format.idx != 0:
                    content_box = ph
                    break
        if not content_box:
            logging.info(f"幻灯片 {slide.slide_id} 无正文占位符，自动添加更大文本框")
            content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5.5))

        tf = content_box.text_frame
        tf.clear()
        tf.word_wrap = True
        paragraphs = content.split('\n\n')
        for i, para in enumerate(paragraphs):
            if para.strip():
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = para.strip()
                p.font.size = Pt(18)
                p.font.name = '微软雅黑'
                p.space_after = Pt(8)
                if layout_type == 'content':
                    p.level = 0
        logging.debug(
            f"幻灯片 {slide.slide_id} 正文字符数: {len(content)}, 段落数: {len([p for p in paragraphs if p.strip()])}")

    def _add_quote(self, slide, content):
        content_box = None
        if len(slide.placeholders) > 1:
            for ph in slide.placeholders:
                if ph.placeholder_format.idx != 0:
                    content_box = ph
                    break
        if not content_box:
            content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5.5))

        tf = content_box.text_frame
        tf.clear()
        tf.text = content
        for p in tf.paragraphs:
            p.font.italic = True
            p.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
            p.alignment = PP_ALIGN.CENTER
        logging.debug(f"引用页正文: {content}")

    def _add_chart(self, slide, content):
        chart_type_match = re.search(r'{chart:(\w+)}', content)
        chart_type = chart_type_match.group(1) if chart_type_match else "default"
        chart_data = ChartData()
        chart_data.categories = ['2021', '2022', '2023']
        chart_data.add_series(chart_type.title(), (120, 180, 250) if chart_type == 'market_size' else (
        15, 25, 30) if chart_type == 'growth_rate' else (100, 200, 300))
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1.5), Inches(10), Inches(5), chart_data
        ).chart
        chart.has_title = True
        chart.chart_title.text_frame.text = f"{chart_type.replace('_', ' ').title()} 数据"
        logging.debug(f"图表页标题: {chart.chart_title.text_frame.text}")

    def _add_end_slide(self):
        self.current_stage = "生成结束页"
        if self.update_status:
            self.update_status(self.current_stage)
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[self.layout_map['end']])
        title = slide.shapes.title or slide.shapes.add_textbox(Inches(1), Inches(1), Inches(11), Inches(1))
        title.text = "感谢聆听"
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x2E, 0x86, 0xC1)

        content_box = None
        if len(slide.placeholders) > 1:
            for ph in slide.placeholders:
                if ph.placeholder_format.idx != 0:
                    content_box = ph
                    break
        if not content_box:
            content_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11), Inches(2))
        content_box.text_frame.text = "欢迎提问\ncontact@example.com"
        logging.debug(f"结束页标题: 感谢聆听")

    def _save_ppt(self, output_path):
        if not self.cancelled:
            self.current_stage = "正在保存PPT"
            if self.update_status:
                self.update_status(self.current_stage)
            output_path = os.path.abspath(output_path)
            logging.info(f"保存PPT时的最终路径: {output_path}")
            self.prs.save(output_path)
            if not os.path.exists(output_path):
                raise FileNotFoundError(f"文件未保存成功: {output_path}")
            logging.info(f"PPT保存成功至: {output_path}")

    def cancel(self):
        self.cancelled = True

class OutlineWorker(QThread):
    outline_generated = pyqtSignal(str)
    error = pyqtSignal(str)

    def __init__(self, generator, content):
        super().__init__()
        self.generator = generator
        self.content = content

    def run(self):
        try:
            if os.path.isfile(self.content):
                with open(self.content, 'r', encoding='utf-8') as f:
                    outline_text = f.read()
            elif re.search(r'^#+\s|\* |- |\d+\.', self.content, re.MULTILINE):
                outline_text = self.content
            else:
                outline_text = self.generator.generate_outline(self.content)
            data = self.generator.parse_markdown(outline_text)
            outline = f"# {data['title']}\n\n"
            for section in data["sections"]:
                outline += f"## {section['title']}\n"
                for content in section["content"]:
                    outline += f"- {content}\n"
                outline += "\n"
            self.outline_generated.emit(outline)
        except Exception as e:
            self.error.emit(f"生成大纲失败: {str(e)}")

class PPTWorker(QThread):
    status_update = pyqtSignal(str)
    finished = pyqtSignal(str)
    error = pyqtSignal(str)

    def __init__(self, generator, outline_text, output_path):
        super().__init__()
        self.generator = generator
        self.outline_text = outline_text
        self.output_path = output_path

    def run(self):
        try:
            data = self.generator.parse_input(self.outline_text)
            path = self.generator.build_ppt(data, self.output_path)
            if path and not self.generator.cancelled:
                self.finished.emit(f"PPT已保存至:\n{path}")
            elif self.generator.cancelled:
                self.finished.emit("PPT生成已被取消")
        except Exception as e:
            self.error.emit(f"发生错误: {str(e)}\n建议: 检查输入内容、模板文件或文件夹权限。")

    def cancel(self):
        self.generator.cancel()

class TranslationWorker(QThread):
    finished = pyqtSignal(str)
    error = pyqtSignal(str)

    def __init__(self, prompt):
        super().__init__()
        self.prompt = prompt

    def run(self):
        try:
            translated_text = call_deepseek_api(self.prompt)
            self.finished.emit(translated_text)
        except Exception as e:
            self.error.emit(str(e))

class PPTGenerationWidget(QWidget):
    def __init__(self, parent=None):
        super().__init__(parent)
        self.running = False
        self.selected_folder = "D:\\outputPPT"
        self.generator = PPTGenerator(self.update_status)
        self.worker = None
        self.outline_worker = None
        self.start_time = None
        self.outline_text = ""
        self.logo_path = "./data_init/logo.png"
        self.init_ui()

    def init_ui(self):
        self.setStyleSheet("""
            background-color: qradialgradient(cx: 0%, cy: 100%, radius: 100%, stop: 0 #00CED1, stop: 1 #87CEEB);
            border-radius: 50px;
            border: 1px solid #AAAAAA;
        """)
        main_layout = QVBoxLayout(self)
        main_layout.setContentsMargins(20, 20, 20, 20)
        main_layout.setSpacing(15)

        input_frame = QFrame(self)
        input_frame.setStyleSheet("background: transparent;")
        input_layout = QHBoxLayout(input_frame)
        self.input_text_radio = QRadioButton("文本输入", input_frame)
        self.input_file_radio = QRadioButton("文件上传", input_frame)
        self.input_text_radio.setChecked(True)
        self.input_text_radio.toggled.connect(self.toggle_input)
        for radio in [self.input_text_radio, self.input_file_radio]:
            radio.setFont(QFont("微软雅黑", 12))
        input_layout.addWidget(self.input_text_radio)
        input_layout.addWidget(self.input_file_radio)
        main_layout.addWidget(input_frame)

        self.text_input = QTextEdit(self)
        self.text_input.setFont(QFont("微软雅黑", 12))
        main_layout.addWidget(self.text_input)

        self.file_frame = QFrame(self)
        self.file_frame.setStyleSheet("background: transparent;")
        file_layout = QHBoxLayout(self.file_frame)
        file_upload_label = QLabel("文件上传:", self.file_frame)
        file_upload_label.setFont(QFont("微软雅黑", 12))
        self.file_path = QLineEdit(self.file_frame)
        self.file_path.setReadOnly(True)
        self.file_path.setFixedHeight(60)
        file_btn = QPushButton("选择文件", self.file_frame)
        file_btn.setFixedHeight(60)
        file_btn.setFont(QFont("微软雅黑", 12))
        file_btn.clicked.connect(self.load_file)
        file_layout.addWidget(file_upload_label)
        file_layout.addWidget(self.file_path)
        file_layout.addWidget(file_btn)
        self.file_frame.hide()
        main_layout.addWidget(self.file_frame)

        template_frame = QFrame(self)
        template_frame.setStyleSheet("background: transparent;")
        template_layout = QHBoxLayout(template_frame)
        template_label = QLabel("选择模板:", template_frame)
        template_label.setFont(QFont("微软雅黑", 12))
        self.template_combo = QComboBox(template_frame)
        self.template_combo.addItems(TEMPLATE_MAP.keys())
        self.template_combo.setCurrentText("简约类")
        self.template_combo.setFixedHeight(60)
        self.template_combo.setFont(QFont("微软雅黑", 12))
        self.template_combo.currentTextChanged.connect(self.update_template)
        self.template_combo.setStyleSheet("""
            QComboBox {
                background-color: rgba(255, 255, 255, 0.9);
                color: #333333;
                border: 1px solid #AAAAAA;
                border-radius: 10px;
                padding: 5px;
            }
            QComboBox::drop-down {border: none; width: 20px;}
            QComboBox::down-arrow {image: url(./data_init/down_arrow.png); width: 10px; height: 10px;}
            QComboBox QAbstractItemView {
                background-color: rgba(255, 255, 255, 0.95);
                color: #333333;
                selection-background-color: rgba(53, 142, 255, 0.8);
                selection-color: white;
                border: 1px solid #AAAAAA;
                border-radius: 5px;
            }
        """)
        template_layout.addWidget(template_label)
        template_layout.addWidget(self.template_combo)
        main_layout.addWidget(template_frame)

        output_frame = QFrame(self)
        output_frame.setStyleSheet("background: transparent;")
        output_layout = QHBoxLayout(output_frame)
        folder_label = QLabel("输出文件夹:", output_frame)
        folder_label.setFont(QFont("微软雅黑", 12))
        self.folder_input = QLineEdit(self.selected_folder, output_frame)
        self.folder_input.setFixedHeight(60)
        folder_btn = QPushButton("选择文件夹", output_frame)
        folder_btn.setFixedHeight(60)
        folder_btn.clicked.connect(self.select_folder)
        folder_btn.setFont(QFont("微软雅黑", 12))
        output_label = QLabel("输出文件:", output_frame)
        output_label.setFont(QFont("微软雅黑", 12))
        self.output_path = QLineEdit(self.get_default_output_path(), output_frame)
        self.output_path.setReadOnly(True)
        self.output_path.setFixedHeight(60)
        output_layout.addWidget(folder_label)
        output_layout.addWidget(self.folder_input)
        output_layout.addWidget(folder_btn)
        output_layout.addWidget(output_label)
        output_layout.addWidget(self.output_path)
        main_layout.addWidget(output_frame)

        self.outline_display = QTextEdit(self)
        self.outline_display.setFont(QFont("微软雅黑", 12))
        main_layout.addWidget(self.outline_display)

        control_frame = QFrame(self)
        control_frame.setStyleSheet("background: transparent;")
        control_layout = QHBoxLayout(control_frame)
        self.progress = QProgressBar(control_frame)
        self.progress.setRange(0, 0)
        self.progress.hide()
        self.preview_btn = QPushButton("生成大纲", control_frame)
        self.preview_btn.setFixedHeight(60)
        self.preview_btn.setFont(QFont("微软雅黑", 12, QFont.Bold))
        self.preview_btn.clicked.connect(self.preview_outline)
        self.generate_btn = QPushButton("确认生成PPT", control_frame)
        self.generate_btn.setFixedHeight(60)
        self.generate_btn.setFont(QFont("微软雅黑", 12, QFont.Bold))
        self.generate_btn.clicked.connect(self.start_generate)
        self.generate_btn.setEnabled(False)
        self.cancel_btn = QPushButton("取消", control_frame)
        self.cancel_btn.setFixedHeight(60)
        self.cancel_btn.setFont(QFont("微软雅黑", 12, QFont.Bold))
        self.cancel_btn.clicked.connect(self.cancel_generate)
        self.cancel_btn.setEnabled(False)
        self.status_label = QLabel("就绪", control_frame)
        self.status_label.setStyleSheet("color: white; background: transparent;")
        self.status_label.setFont(QFont("微软雅黑", 12))
        control_layout.addWidget(self.progress)
        control_layout.addWidget(self.preview_btn)
        control_layout.addWidget(self.generate_btn)
        control_layout.addWidget(self.cancel_btn)
        control_layout.addWidget(self.status_label)
        main_layout.addWidget(control_frame)

        for btn in [file_btn, folder_btn, self.preview_btn, self.generate_btn, self.cancel_btn]:
            btn.setStyleSheet("""
                QPushButton {
                    background: rgba(53, 142, 255, 0.9);
                    color: white;
                    border: none;
                    border-radius: 10px;
                    padding: 10px 20px;
                }
                QPushButton:disabled {background: rgba(150, 150, 150, 0.6);}
                QPushButton:hover {background: rgba(53, 142, 255, 1);}
            """)

    def toggle_input(self):
        if self.input_text_radio.isChecked():
            self.text_input.show()
            self.file_frame.hide()
        else:
            self.text_input.hide()
            self.file_frame.show()

    def load_file(self):
        path, _ = QFileDialog.getOpenFileName(self, "选择文件", "", "Markdown文件 (*.md)")
        if path:
            self.file_path.setText(path)
            with open(path, 'r', encoding='utf-8') as f:
                self.text_input.setPlainText(f.read())

    def load_file_content(self):
        try:
            with open(self.file_path.text(), 'r', encoding='utf-8') as f:
                return f.read()
        except Exception as e:
            logging.error(f"加载文件失败: {str(e)}")
            self.status_label.setText(f"错误: 无法加载文件")
            return ""

    def on_outline_generated(self, outline_text):
        self.outline_display.setPlainText(outline_text)
        self.generate_btn.setEnabled(True)
        self.running = False
        self.progress.hide()
        self.status_label.setText("就绪")
        self.preview_btn.setEnabled(True)

    def on_outline_error(self, error_message):
        logging.error(error_message)
        self.status_label.setText(error_message)
        self.running = False
        self.progress.hide()
        self.preview_btn.setEnabled(True)
        self.generate_btn.setEnabled(False)
        QMessageBox.critical(self, "错误", error_message)

    def select_folder(self):
        folder = QFileDialog.getExistingDirectory(self, "选择文件夹", self.selected_folder)
        if folder:
            self.selected_folder = folder
            self.folder_input.setText(folder)
            self.output_path.setText(self.get_default_output_path())

    def get_default_output_path(self):
        folder = os.path.abspath(self.selected_folder)
        filename = datetime.now().strftime("%Y%m%d%H%M") + ".pptx"
        output_path = os.path.join(folder, filename)
        os.makedirs(folder, exist_ok=True)
        return output_path

    def update_template(self, template_name):
        self.generator.set_template(template_name)

    def preview_outline(self):
        if not self.running:
            self.running = True
            self.progress.show()
            self.status_label.setText("生成大纲中...")
            self.preview_btn.setEnabled(False)
            self.generate_btn.setEnabled(False)
            content = self.text_input.toPlainText() if self.input_text_radio.isChecked() else self.load_file_content()
            self.outline_worker = OutlineWorker(self.generator, content)
            self.outline_worker.outline_generated.connect(self.on_outline_generated)
            self.outline_worker.error.connect(self.on_outline_error)
            self.outline_worker.start()

    def start_generate(self):
        if not self.running:
            self.running = True
            self.start_time = datetime.now()
            self.progress.show()
            self.status_label.setText("生成PPT中...")
            self.generate_btn.setEnabled(False)
            self.cancel_btn.setEnabled(True)
            self.output_path.setText(self.get_default_output_path())
            self.outline_text = self.outline_display.toPlainText()
            self.worker = PPTWorker(self.generator, self.outline_text, self.output_path.text())
            self.worker.status_update.connect(self.update_status)
            self.worker.finished.connect(self.on_generate_finished)
            self.worker.error.connect(self.on_generate_error)
            self.worker.start()
            self.update_ui()

    def cancel_generate(self):
        if self.running and self.worker:
            self.worker.cancel()
            self.status_label.setText("取消中...")

    def on_generate_finished(self, message):
        self.running = False
        self.progress.hide()
        self.status_label.setText("就绪")
        self.preview_btn.setEnabled(True)
        self.generate_btn.setEnabled(True)
        self.cancel_btn.setEnabled(False)
        QMessageBox.information(self, "成功", message)

    def on_generate_error(self, message):
        self.running = False
        self.progress.hide()
        self.status_label.setText("就绪")
        self.preview_btn.setEnabled(True)
        self.generate_btn.setEnabled(True)
        self.cancel_btn.setEnabled(False)
        QMessageBox.critical(self, "错误", message)

    def update_status(self, stage):
        self.generator.current_stage = stage

    def update_ui(self):
        if self.running and self.start_time:
            elapsed = (datetime.now() - self.start_time).total_seconds()
            self.status_label.setText(f"{self.generator.current_stage} (耗时: {elapsed:.2f}秒)")
            QTimer.singleShot(100, self.update_ui)
        elif not self.running:
            self.status_label.setText("就绪")

class TranslationWidget(QWidget):
    def __init__(self, parent=None):
        super().__init__(parent)
        self.init_ui()

    def init_ui(self):
        self.setStyleSheet("""
            background-color: qradialgradient(cx: 0%, cy: 100%, radius: 100%, stop: 0 #00CED1, stop: 1 #87CEEB);
            border-radius: 50px;
            border: 1px solid #AAAAAA;
        """)
        main_layout = QVBoxLayout(self)
        main_layout.setContentsMargins(20, 20, 20, 20)
        main_layout.setSpacing(15)

        # Direction selection
        direction_frame = QFrame(self)
        direction_frame.setStyleSheet("background: transparent;")
        direction_layout = QHBoxLayout(direction_frame)
        self.zh_to_en_radio = QRadioButton("中文 -> 英文", direction_frame)
        self.en_to_zh_radio = QRadioButton("英文 -> 中文", direction_frame)
        self.zh_to_en_radio.setChecked(True)
        for radio in [self.zh_to_en_radio, self.en_to_zh_radio]:
            radio.setFont(QFont("微软雅黑", 12))
        direction_layout.addWidget(self.zh_to_en_radio)
        direction_layout.addWidget(self.en_to_zh_radio)
        main_layout.addWidget(direction_frame)

        # Input text
        input_label = QLabel("输入文本:", self)
        input_label.setFont(QFont("微软雅黑", 12))
        main_layout.addWidget(input_label)
        self.input_text = QTextEdit(self)
        self.input_text.setFont(QFont("微软雅黑", 12))
        main_layout.addWidget(self.input_text)

        # Translate button
        self.translate_btn = QPushButton("翻译", self)
        self.translate_btn.setFixedHeight(60)
        self.translate_btn.setFont(QFont("微软雅黑", 12, QFont.Bold))
        self.translate_btn.clicked.connect(self.start_translation)
        self.translate_btn.setStyleSheet("""
            QPushButton {
                background: rgba(53, 142, 255, 0.9);
                color: white;
                border: none;
                border-radius: 10px;
                padding: 10px 20px;
            }
            QPushButton:disabled {background: rgba(150, 150, 150, 0.6);}
            QPushButton:hover {background: rgba(53, 142, 255, 1);}
        """)
        main_layout.addWidget(self.translate_btn)

        # Output text
        output_label = QLabel("翻译结果:", self)
        output_label.setFont(QFont("微软雅黑", 12))
        main_layout.addWidget(output_label)
        self.output_text = QTextEdit(self)
        self.output_text.setFont(QFont("微软雅黑", 12))
        self.output_text.setReadOnly(True)
        main_layout.addWidget(self.output_text)

        # Status label
        self.status_label = QLabel("就绪", self)
        self.status_label.setStyleSheet("color: white; background: transparent;")
        self.status_label.setFont(QFont("微软雅黑", 12))
        main_layout.addWidget(self.status_label)

    def start_translation(self):
        text = self.input_text.toPlainText().strip()
        if not text:
            QMessageBox.warning(self, "警告", "请输入要翻译的文本")
            return
        if self.zh_to_en_radio.isChecked():
            prompt = f"Translate the following Chinese text to English:\n\n{text}"
        else:
            prompt = f"Translate the following English text to Chinese:\n\n{text}"
        self.translate_btn.setEnabled(False)
        self.status_label.setText("翻译中...")
        self.worker = TranslationWorker(prompt)
        self.worker.finished.connect(self.on_translation_finished)
        self.worker.error.connect(self.on_translation_error)
        self.worker.start()

    def on_translation_finished(self, translated_text):
        self.output_text.setPlainText(translated_text)
        self.translate_btn.setEnabled(True)
        self.status_label.setText("完成")

    def on_translation_error(self, error_message):
        QMessageBox.critical(self, "错误", f"翻译失败: {error_message}")
        self.translate_btn.setEnabled(True)
        self.status_label.setText("就绪")

class IntegratedApp(QMainWindow):
    def __init__(self):
        super().__init__()
        self.setWindowFlags(Qt.FramelessWindowHint)
        self.logo_path = "./data_init/logo.png"
        self.background_path = "./data_init/background.jpg"
        self.chart_worker = None  # 用于跟踪图表生成线程

        # 初始化时间标签
        self.time_label = QLabel()
        self.time_label.setStyleSheet("color: white; background: transparent; border: none; font-size: 20px;")
        self.time_label.setFont(QFont("微软雅黑", 12))
        self.update_time()  # 初始更新时间

        # 设置定时器，每秒更新时间
        self.timer = QTimer(self)
        self.timer.timeout.connect(self.update_time)
        self.timer.start(1000)  # 每1000毫秒（1秒）更新一次



        self.init_ui()

    def update_time(self):
        current_time = QDateTime.currentDateTime().toString("yyyy-MM-dd HH:mm:ss")
        self.time_label.setText(current_time)

    def init_ui(self):
        self.setWindowTitle("智能生成工具")
        self.setGeometry(100, 100, 1200, 900)
        logging.info("初始化窗口大小: 1200x900")

        central_widget = QWidget(self)
        central_widget.setStyleSheet("""
            background-color: qradialgradient(cx: 0%, cy: 100%, radius: 100%, stop: 0 #00CED1, stop: 1 #87CEEB);
            border-radius: 50px;
            border: 0px solid #AAAAAA;
        """)
        self.setCentralWidget(central_widget)

        main_layout = QVBoxLayout(central_widget)
        main_layout.setContentsMargins(20, 20, 20, 20)
        main_layout.setSpacing(0)

        top_layout = QHBoxLayout()

        left_panel = QWidget()
        left_panel.setStyleSheet("background: rgba(255, 255, 255, 0.9); border-radius: 10px;")
        left_panel_layout = QVBoxLayout(left_panel)
        left_panel_layout.setAlignment(Qt.AlignTop)
        left_panel_layout.setContentsMargins(10, 10, 10, 10)

        self.logo_label = QLabel()
        button_width = 180
        self.logo_label.setFixedSize(button_width, button_width)
        self.logo_label.setStyleSheet("background: transparent;")  # 添加透明背景
        self.update_logo()
        left_panel_layout.addWidget(self.logo_label)

        self.generate_btn = QPushButton("智能生成工具")
        self.generate_btn.setFixedHeight(80)
        self.generate_btn.clicked.connect(self.show_generate_menu)
        self.generate_btn.setStyleSheet("""
            QPushButton {background: rgba(53, 142, 255, 0.9); color: white; border-radius: 5px; padding: 10px;}
            QPushButton:hover {background: rgba(53, 142, 255, 1);}
        """)
        left_panel_layout.addWidget(self.generate_btn)

        self.help_btn = QPushButton("帮助")
        self.help_btn.setFixedHeight(80)
        self.help_btn.clicked.connect(self.show_help)
        self.help_btn.setStyleSheet("""
            QPushButton {background: rgba(53, 142, 255, 0.9); color: white; border-radius: 5px; padding: 10px;}
            QPushButton:hover {background: rgba(53, 142, 255, 1);}
        """)
        left_panel_layout.addWidget(self.help_btn)

        self.settings_btn = QPushButton("设置")
        self.settings_btn.setFixedHeight(80)
        self.settings_btn.clicked.connect(self.set_background_image)
        self.settings_btn.setStyleSheet("""
            QPushButton {background: rgba(53, 142, 255, 0.9); color: white; border-radius: 5px; padding: 10px;}
            QPushButton:hover {background: rgba(53, 142, 255, 1);}
        """)
        left_panel_layout.addWidget(self.settings_btn)

        left_panel_layout.addStretch()
        left_panel.setFixedWidth(200)
        left_panel.setSizePolicy(QSizePolicy.Fixed, QSizePolicy.Expanding)
        top_layout.addWidget(left_panel)

        right_container = QWidget()
        right_layout = QVBoxLayout(right_container)
        right_layout.setContentsMargins(10, 0, 0, 0)

        title_control_layout = QHBoxLayout()
        self.title_label = QLabel("智能生成工具")
        self.title_label.setStyleSheet("color: #333333; background: transparent;")
        self.title_label.setFixedHeight(50)
        font = QFont()
        font.setPointSize(20)
        self.title_label.setFont(font)

        minimize_btn = QPushButton("—")
        minimize_btn.clicked.connect(self.showMinimized)
        close_btn = QPushButton("×")
        close_btn.clicked.connect(self.close)
        for btn in [minimize_btn, close_btn]:
            btn.setFixedSize(30, 30)
            btn.setStyleSheet("""
                QPushButton {background: transparent; color: #666666; border: none;}
                QPushButton:hover {color: #FF0000;}
            """)

        title_control_layout.addWidget(self.title_label)
        title_control_layout.addStretch()
        title_control_layout.addWidget(minimize_btn)
        title_control_layout.addWidget(close_btn)
        right_layout.addLayout(title_control_layout)



        self.stack = QStackedWidget()
        self.stack.addWidget(self.create_welcome_widget())
        self.stack.addWidget(PPTGenerationWidget(self))
        self.stack.addWidget(self.create_chart_widget())
        self.stack.addWidget(self.create_help())
        self.stack.addWidget(self.create_selection_widget())
        self.stack.addWidget(TranslationWidget(self))
        right_layout.addWidget(self.stack)

        top_layout.addWidget(right_container)
        top_layout.setStretch(0, 1)
        top_layout.setStretch(1, 3)
        main_layout.addLayout(top_layout)

        self.update_title_color()
        logging.info("界面初始化完成")

    def get_default_chart_output_path(self):
        folder = self.chart_folder_input.text()
        format_str = self.format_combo.currentText().lower()
        filename = f"chart.{format_str}"
        return os.path.join(folder, filename)

    def load_chart_file(self):
        file_path, _ = QFileDialog.getOpenFileName(self, "选择Excel文件", "",
                                                   "Excel Files (*.xlsx *.xls);;All Files (*)")
        if file_path:
            self.chart_file_path.setText(file_path)

    def select_chart_folder(self):
        folder = QFileDialog.getExistingDirectory(self, "选择输出文件夹", self.chart_folder_input.text())
        if folder:
            self.chart_folder_input.setText(folder)
            self.chart_output_path.setText(self.get_default_chart_output_path())

    def update_chart_output_path_extension(self):
        self.chart_output_path.setText(self.get_default_chart_output_path())

    def start_chart_generate(self):
        file_path = self.chart_file_path.text().strip()
        chart_type = self.chart_type_combo.currentText()
        output_path = self.chart_output_path.text().strip()

        if not file_path or not output_path:
            QMessageBox.warning(self, "警告", "请选择Excel文件和输出路径")
            return

        if not os.path.exists(file_path):
            QMessageBox.warning(self, "警告", "Excel文件不存在")
            return

        self.chart_worker = ChartWorker(file_path, chart_type, output_path)
        self.chart_worker.status_update.connect(self.update_chart_status)
        self.chart_worker.finished.connect(self.on_chart_finished)
        self.chart_worker.error.connect(self.on_chart_error)
        self.chart_worker.start()

        self.chart_progress.show()
        self.chart_status_label.setText("生成图表中...")
        self.chart_generate_btn.setEnabled(False)
        self.chart_cancel_btn.setEnabled(True)

    def cancel_chart_generate(self):
        if self.chart_worker is not None:
            self.chart_worker.cancel()
            self.chart_status_label.setText("取消中...")

    def update_chart_status(self, message):
        self.chart_status_label.setText(message)

    def on_chart_finished(self, message):
        self.chart_progress.hide()
        self.chart_status_label.setText("就绪")
        self.chart_generate_btn.setEnabled(True)
        self.chart_cancel_btn.setEnabled(False)
        QMessageBox.information(self, "成功", message)

    def on_chart_error(self, message):
        self.chart_progress.hide()
        self.chart_status_label.setText("就绪")
        self.chart_generate_btn.setEnabled(True)
        self.chart_cancel_btn.setEnabled(False)
        QMessageBox.critical(self, "错误", message)

    def create_welcome_widget(self):
        widget = QWidget()
        layout = QVBoxLayout(widget)
        self.bg_label = QLabel()
        pixmap = QPixmap(self.background_path).scaled(self.size(), Qt.KeepAspectRatio, Qt.SmoothTransformation)
        self.bg_label.setPixmap(pixmap)
        self.bg_label.setAlignment(Qt.AlignCenter)

        # 添加时间标签到布局
        layout.addWidget(self.bg_label)
        layout.addStretch()  # 推到下方
        time_layout = QHBoxLayout()
        time_layout.addStretch()  # 推到右侧
        time_layout.addWidget(self.time_label)
        layout.addLayout(time_layout)

        return widget

    def create_chart_widget(self):
        widget = QWidget()
        layout = QVBoxLayout(widget)
        label_font = QFont()
        label_font.setPointSize(14)

        # 文件选择框架
        file_frame = QFrame()
        file_frame.setStyleSheet("background: transparent;")
        file_layout = QHBoxLayout(file_frame)
        file_label = QLabel("选择 Excel 文件:", file_frame)
        file_label.setFont(QFont("微软雅黑", 12))
        self.chart_file_path = QLineEdit(file_frame)
        self.chart_file_path.setReadOnly(True)
        self.chart_file_path.setFixedHeight(60)
        file_btn = QPushButton("选择文件", file_frame)
        file_btn.setFixedHeight(60)
        file_btn.setFont(QFont("微软雅黑", 12))
        file_btn.clicked.connect(self.load_chart_file)
        file_layout.addWidget(file_label)
        file_layout.addWidget(self.chart_file_path)
        file_layout.addWidget(file_btn)
        layout.addWidget(file_frame)

        # 图表类型选择框架
        chart_frame = QFrame()
        chart_frame.setStyleSheet("background: transparent;")
        chart_layout = QHBoxLayout(chart_frame)
        chart_label = QLabel("图表类型:", chart_frame)
        chart_label.setFont(QFont("微软雅黑", 12))
        self.chart_type_combo = QComboBox(chart_frame)
        self.chart_type_combo.addItems(["饼状图", "柱状图", "折线图", "散点图", "热力图", "直方图", "气泡图"])
        self.chart_type_combo.setFixedHeight(60)
        self.chart_type_combo.setFont(QFont("微软雅黑", 12))
        self.chart_type_combo.setStyleSheet("""
            QComboBox {
                background-color: rgba(255, 255, 255, 0.9);
                color: #333333;
                border: 1px solid #AAAAAA;
                border-radius: 10px;
                padding: 5px;
            }
            QComboBox::drop-down {border: none; width: 20px;}
            QComboBox::down-arrow {image: url(./data_init/down_arrow.png); width: 10px; height: 10px;}
            QComboBox QAbstractItemView {
                background-color: rgba(255, 255, 255, 0.95);
                color: #333333;
                selection-background-color: rgba(53, 142, 255, 0.8);
                selection-color: white;
                border: 1px solid #AAAAAA;
                border-radius: 5px;
            }
        """)
        chart_layout.addWidget(chart_label)
        chart_layout.addWidget(self.chart_type_combo)
        layout.addWidget(chart_frame)

        # 输出格式选择框架
        format_frame = QFrame()
        format_frame.setStyleSheet("background: transparent;")
        format_layout = QHBoxLayout(format_frame)
        format_label = QLabel("输出格式:", format_frame)
        format_label.setFont(QFont("微软雅黑", 12))
        self.format_combo = QComboBox(format_frame)
        self.format_combo.addItems(["PNG", "PDF", "SVG"])
        self.format_combo.setFixedHeight(60)
        self.format_combo.setFont(QFont("微软雅黑", 12))
        self.format_combo.currentTextChanged.connect(self.update_chart_output_path_extension)
        self.format_combo.setStyleSheet("""
            QComboBox {
                background-color: rgba(255, 255, 255, 0.9);
                color: #333333;
                border: 1px solid #AAAAAA;
                border-radius: 10px;
                padding: 5px;
            }
            QComboBox::drop-down {border: none; width: 20px;}
            QComboBox::down-arrow {image: url(./data_init/down_arrow.png); width: 10px; height: 10px;}
            QComboBox QAbstractItemView {
                background-color: rgba(255, 255, 255, 0.95);
                color: #333333;
                selection-background-color: rgba(53, 142, 255, 0.8);
                selection-color: white;
                border: 1px solid #AAAAAA;
                border-radius: 5px;
            }
        """)
        format_layout.addWidget(format_label)
        format_layout.addWidget(self.format_combo)
        layout.addWidget(format_frame)

        # 输出路径框架
        output_frame = QFrame()
        output_frame.setStyleSheet("background: transparent;")
        output_layout = QHBoxLayout(output_frame)
        folder_label = QLabel("输出文件夹:", output_frame)
        folder_label.setFont(QFont("微软雅黑", 12))
        self.chart_folder_input = QLineEdit("D:/outputCharts", output_frame)
        self.chart_folder_input.setFixedHeight(60)
        folder_btn = QPushButton("选择文件夹", output_frame)
        folder_btn.setFixedHeight(60)
        folder_btn.setFont(QFont("微软雅黑", 12))
        folder_btn.clicked.connect(self.select_chart_folder)
        output_label = QLabel("输出文件:", output_frame)
        output_label.setFont(QFont("微软雅黑", 12))
        self.chart_output_path = QLineEdit(self.get_default_chart_output_path(), output_frame)
        self.chart_output_path.setReadOnly(True)
        self.chart_output_path.setFixedHeight(60)
        output_layout.addWidget(folder_label)
        output_layout.addWidget(self.chart_folder_input)
        output_layout.addWidget(folder_btn)
        output_layout.addWidget(output_label)
        output_layout.addWidget(self.chart_output_path)
        layout.addWidget(output_frame)

        # 控制按钮框架
        control_frame = QFrame()
        control_frame.setStyleSheet("background: transparent;")
        control_layout = QHBoxLayout(control_frame)
        self.chart_progress = QProgressBar(control_frame)
        self.chart_progress.setRange(0, 0)
        self.chart_progress.hide()
        self.chart_generate_btn = QPushButton("生成图表", control_frame)
        self.chart_generate_btn.setFixedHeight(60)
        self.chart_generate_btn.setFont(QFont("微软雅黑", 12, QFont.Bold))
        self.chart_generate_btn.clicked.connect(self.start_chart_generate)
        self.chart_cancel_btn = QPushButton("取消", control_frame)
        self.chart_cancel_btn.setFixedHeight(60)
        self.chart_cancel_btn.setFont(QFont("微软雅黑", 12, QFont.Bold))
        self.chart_cancel_btn.clicked.connect(self.cancel_chart_generate)
        self.chart_cancel_btn.setEnabled(False)
        self.chart_status_label = QLabel("就绪", control_frame)
        self.chart_status_label.setStyleSheet("color: white; background: transparent;")
        self.chart_status_label.setFont(QFont("微软雅黑", 12))
        control_layout.addWidget(self.chart_progress)
        control_layout.addWidget(self.chart_generate_btn)
        control_layout.addWidget(self.chart_cancel_btn)
        control_layout.addWidget(self.chart_status_label)
        layout.addWidget(control_frame)

        # 统一按钮样式
        for btn in [file_btn, folder_btn, self.chart_generate_btn, self.chart_cancel_btn]:
            btn.setStyleSheet("""
                QPushButton {
                    background: rgba(53, 142, 255, 0.9); 
                    color: white; 
                    border: none; 
                    border-radius: 10px; 
                    padding: 10px 20px;
                }
                QPushButton:disabled {background: rgba(150, 150, 150, 0.6);}
                QPushButton:hover {background: rgba(53, 142, 255, 1);}
            """)

        return widget

    def create_help(self):
        widget = QWidget()
        layout = QVBoxLayout(widget)
        help_text = QTextBrowser()
        help_text.setText("""
    ### 智能生成工具 - 使用说明

    欢迎使用智能生成工具！本工具旨在帮助您快速生成PPT演示文稿和数据图表，提升工作效率。以下是功能介绍和操作指南：

    #### 一、功能概览
    1. **PPT生成**  
       根据输入的文本内容或上传的文件，自动生成结构化的PPT演示文稿。支持多种模板选择，适用于教学、政务、学术等场景。
    2. **图表生成**  
       从Excel文件导入数据，生成多种类型的图表（如饼状图、柱状图、热力图等），并支持多种输出格式（PNG、PDF、SVG）。
    3. **设置**  
       通过左侧“设置”按钮，用户可自定义欢迎界面的背景图片。

    #### 二、操作指南

    ##### 1. 生成PPT
    - **步骤：**
      1. **选择输入方式**  
         - 在“文本输入”模式下，直接在文本框中输入PPT的大纲或内容。  
         - 在“文件上传”模式下，点击“选择文件”按钮上传包含大纲的文本文件（支持*.txt格式）。
      2. **选择模板**  
         - 在“选择模板”下拉菜单中，选择适合的模板（如“教学课件类”、“政务报告类”）。默认使用“简约类”模板。
      3. **设置输出路径**  
         - 点击“选择文件夹”指定输出目录，默认路径为“D:/outputPPT/output.pptx”。  
         - 输出文件名可根据需要调整。
      4. **生成大纲**  
         - 点击“生成大纲”按钮，系统将解析输入内容并生成结构化大纲，显示在下方文本框中。
      5. **确认生成PPT**  
         - 检查大纲无误后，点击“确认生成PPT”按钮开始生成。生成过程中可通过“取消”按钮中止操作。
    - **提示：**  
      - 输入内容应简洁清晰，避免使用违法或不当语言。  
      - 生成状态将在界面底部显示，完成后会弹出成功提示。

    ##### 2. 生成图表
    - **步骤：**
      1. **选择Excel文件**  
         - 点击“选择文件”按钮，上传包含数据的Excel文件（支持*.xlsx和*.xls格式）。  
         - 确保文件中数据格式规范，第一列通常为X轴数据，第二列为Y轴数据（部分图表可能需要更多列）。
      2. **选择图表类型**  
         - 在“图表类型”下拉菜单中选择所需类型（如“饼状图”、“折线图”）。不同图表适用于不同数据展示需求。
      3. **选择输出格式**  
         - 在“输出格式”下拉菜单中选择文件格式（PNG、PDF、SVG），默认为PNG。
      4. **设置输出路径**  
         - 点击“选择文件夹”指定输出目录，默认路径为“D:/outputCharts/chart.[格式]”。
      5. **生成图表**  
         - 点击“生成图表”按钮开始生成。生成过程中可通过“取消”按钮中止操作。
    - **提示：**  
      - Excel文件需包含有效数据，空文件或格式错误将导致生成失败。  
      - 生成的图表将自动保存至指定路径，完成后会弹出成功提示。

    ##### 3. 设置背景图片
    - **步骤：**
      1. 点击左侧选项栏的“设置”按钮。
      2. 在弹出的文件选择对话框中，选择一张图片（支持*.jpg、*.png、*.bmp格式）。
      3. 确认后，欢迎界面的背景将更新为所选图片。
    - **提示：**  
      - 建议选择与窗口大小（1200x900）匹配的图片，以获得最佳显示效果。

    ##### 4. 界面导航
    - **智能生成工具按钮**：返回功能选择页面，可切换至PPT或图表生成模式。
    - **帮助按钮**：查看本使用说明。
    - **设置按钮**：自定义背景图片。
    - **窗口操作**：支持拖动窗口、点击“—”最小化、“×”关闭程序。

    #### 三、注意事项
    1. **合法合规**  
       - 请确保输入内容和生成的文件符合中国法律法规及道德规范，不得包含违法、暴力或敏感政治内容。
    2. **文件权限**  
       - 确保输出文件夹具有写入权限，否则可能导致保存失败。
    3. **系统要求**  
       - 本工具依赖Python环境及相关库（如PyQt5、pandas、matplotlib等），请确保运行环境正常。
    4. **日志记录**  
       - 操作日志将保存至程序运行目录下的“app_YYYYMMDD.log”文件，便于排查问题。

    #### 四、常见问题
    - **Q：生成失败怎么办？**  
      - A：检查输入内容是否为空、文件路径是否有效，或查看错误提示调整操作。
    - **Q：支持哪些图表类型？**  
      - A：支持饼状图、柱状图、折线图、散点图、热力图、直方图和气泡图。
    - **Q：如何调整PPT内容？**  
      - A：可在生成大纲后手动编辑文本框中的内容，再确认生成。

    如需进一步帮助，请联系技术支持。祝您使用愉快！
        """)
        help_text.setFont(QFont("微软雅黑", 12))
        layout.addWidget(help_text)
        return widget

    def create_selection_widget(self):
        widget = QWidget()
        main_layout = QVBoxLayout(widget)
        main_layout.setContentsMargins(20, 20, 20, 20)
        main_layout.setSpacing(15)

        desc_frame = QFrame(widget)
        desc_frame.setStyleSheet("background: transparent;")
        desc_layout = QHBoxLayout(desc_frame)
        desc_label = QLabel("本系统支持根据描述生成PPT，或将Excel表格导入后进行可视化生成图表", desc_frame)
        desc_label.setFont(QFont("微软雅黑", 14))
        desc_label.setStyleSheet("""
            color: #333333;
            background: transparent;
            border: 1px solid #AAAAAA;
            border-radius: 5px;
            padding: 10px;
        """)
        desc_label.setAlignment(Qt.AlignCenter)
        desc_label.setWordWrap(True)
        desc_layout.addWidget(desc_label)
        main_layout.addWidget(desc_frame)

        btn_frame = QFrame(widget)
        btn_frame.setStyleSheet("background: transparent;")
        btn_layout = QVBoxLayout(btn_frame)
        btn_layout.setSpacing(20)

        # First row: PPT and Chart buttons
        first_row = QHBoxLayout()
        first_row.setSpacing(20)

        ppt_container = QWidget()
        ppt_container.setStyleSheet("background: transparent; border: none;")
        ppt_layout = QVBoxLayout(ppt_container)
        ppt_layout.setAlignment(Qt.AlignCenter)
        ppt_layout.setSpacing(10)

        self.ppt_btn = QPushButton(self)
        self.ppt_btn.setFixedSize(590, 390)
        ppt_bg_path = os.path.join(os.path.dirname(__file__), "data_init", "ppt_button_bg.jpg")
        if not os.path.exists(ppt_bg_path):
            logging.error(f"PPT按钮背景图片不存在: {ppt_bg_path}")
        else:
            ppt_pixmap = QPixmap(ppt_bg_path)
            if not ppt_pixmap.isNull():
                self.ppt_btn.setIcon(QIcon(ppt_pixmap))
                self.ppt_btn.setIconSize(QSize(590, 390))
        ppt_bg_path_fixed = ppt_bg_path.replace('\\', '/')
        ppt_bg_url = f"file:///{ppt_bg_path_fixed}"
        self.ppt_btn.setStyleSheet(
            "QPushButton {"
            f"background-image: url({ppt_bg_url});"
            "background-position: center;"
            "background-repeat: no-repeat;"
            "background-size: contain;"
            "border: none;"
            "border-radius: 10px;"
            "}"
            "QPushButton:hover {"
            f"background-image: url({ppt_bg_url});"
            "background-color: rgba(50, 50, 50, 0.5);"
            "}"
        )
        self.ppt_btn.clicked.connect(lambda: self.stack.setCurrentIndex(1))

        ppt_label = QLabel("生成 PPT", self)
        ppt_label.setFont(QFont("微软雅黑", 16, QFont.Bold))
        ppt_label.setStyleSheet("color: #000000; background: transparent; border: none;")
        ppt_label.setAlignment(Qt.AlignCenter)

        ppt_layout.addWidget(self.ppt_btn)
        ppt_layout.addWidget(ppt_label)

        chart_container = QWidget()
        chart_container.setStyleSheet("background: transparent; border: none;")
        chart_layout = QVBoxLayout(chart_container)
        chart_layout.setAlignment(Qt.AlignCenter)
        chart_layout.setSpacing(10)

        self.chart_btn = QPushButton(self)
        self.chart_btn.setFixedSize(590, 390)
        chart_bg_path = os.path.join(os.path.dirname(__file__), "data_init", "chart_button_bg.jpg")
        if not os.path.exists(chart_bg_path):
            logging.error(f"Chart按钮背景图片不存在: {chart_bg_path}")
        else:
            chart_pixmap = QPixmap(chart_bg_path)
            if not chart_pixmap.isNull():
                self.chart_btn.setIcon(QIcon(chart_pixmap))
                self.chart_btn.setIconSize(QSize(590, 390))
        chart_bg_path_fixed = chart_bg_path.replace('\\', '/')
        chart_bg_url = f"file:///{chart_bg_path_fixed}"
        self.chart_btn.setStyleSheet(
            "QPushButton {"
            f"background-image: url({chart_bg_url});"
            "background-position: center;"
            "background-repeat: no-repeat;"
            "background-size: contain;"
            "border: none;"
            "border-radius: 10px;"
            "}"
            "QPushButton:hover {"
            f"background-image: url({chart_bg_url});"
            "background-color: rgba(50, 50, 50, 0.5);"
            "}"
        )
        self.chart_btn.clicked.connect(lambda: self.stack.setCurrentIndex(2))

        chart_label = QLabel("生成图表", self)
        chart_label.setFont(QFont("微软雅黑", 16, QFont.Bold))
        chart_label.setStyleSheet("color: #000000; background: transparent; border: none;")
        chart_label.setAlignment(Qt.AlignCenter)

        chart_layout.addWidget(self.chart_btn)
        chart_layout.addWidget(chart_label)

        first_row.addWidget(ppt_container)
        first_row.addWidget(chart_container)
        btn_layout.addLayout(first_row)

        # Second row: Translation button (centered)
        second_row = QHBoxLayout()
        second_row.addStretch()

        translation_container = QWidget()
        translation_container.setStyleSheet("background: transparent; border: none;")
        translation_layout = QVBoxLayout(translation_container)
        translation_layout.setAlignment(Qt.AlignCenter)
        translation_layout.setSpacing(10)

        self.translation_btn = QPushButton(self)
        self.translation_btn.setFixedSize(590, 390)
        translation_bg_path = os.path.join(os.path.dirname(__file__), "data_init", "translation_button_bg.jpg")
        if not os.path.exists(translation_bg_path):
            logging.error(f"Translation按钮背景图片不存在: {translation_bg_path}")
        else:
            translation_pixmap = QPixmap(translation_bg_path)
            if not translation_pixmap.isNull():
                self.translation_btn.setIcon(QIcon(translation_pixmap))
                self.translation_btn.setIconSize(QSize(590, 390))
        translation_bg_path_fixed = translation_bg_path.replace('\\', '/')
        translation_bg_url = f"file:///{translation_bg_path_fixed}"
        self.translation_btn.setStyleSheet(
            "QPushButton {"
            f"background-image: url({translation_bg_url});"
            "background-position: center;"
            "background-repeat: no-repeat;"
            "background-size: contain;"
            "border: none;"
            "border-radius: 10px;"
            "}"
            "QPushButton:hover {"
            f"background-image: url({translation_bg_url});"
            "background-color: rgba(50, 50, 50, 0.5);"
            "}"
        )
        self.translation_btn.clicked.connect(lambda: self.stack.setCurrentIndex(5))

        translation_label = QLabel("中英互译", self)
        translation_label.setFont(QFont("微软雅黑", 16, QFont.Bold))
        translation_label.setStyleSheet("color: #000000; background: transparent; border: none;")
        translation_label.setAlignment(Qt.AlignCenter)

        translation_layout.addWidget(self.translation_btn)
        translation_layout.addWidget(translation_label)

        second_row.addWidget(translation_container)
        second_row.addStretch()
        btn_layout.addLayout(second_row)

        main_layout.addWidget(btn_frame)

        main_layout.addStretch()
        time_layout = QHBoxLayout()
        time_layout.addStretch()
        time_layout.addWidget(self.time_label)
        main_layout.addLayout(time_layout)

        return widget

    def set_background_image(self):
        file_path, _ = QFileDialog.getOpenFileName(self, "选择背景图片", "", "Image Files (*.jpg *.png *.bmp);;All Files (*)")
        if file_path:
            self.background_path = file_path
            pixmap = QPixmap(self.background_path).scaled(self.size(), Qt.KeepAspectRatio, Qt.SmoothTransformation)
            self.bg_label.setPixmap(pixmap)
            logging.info(f"背景图片已设置为: {self.background_path}")
            QMessageBox.information(self, "成功", "背景图片已更新！")

    def show_generate_menu(self):
        self.stack.setCurrentIndex(4)

    def show_help(self):
        self.stack.setCurrentIndex(3)

    def update_logo(self):
        if os.path.exists(self.logo_path):
            pixmap = QPixmap(self.logo_path).scaled(self.logo_label.size(), Qt.KeepAspectRatio, Qt.SmoothTransformation)
            self.logo_label.setPixmap(pixmap)
        else:
            logging.warning(f"Logo 文件不存在: {self.logo_path}")

    def update_title_color(self):
        r, g, b = [random.randint(0, 255) for _ in range(3)]
        self.title_label.setStyleSheet(f"color: rgb({r}, {g}, {b}); background: transparent;")
        QTimer.singleShot(500, self.update_title_color)

    def mousePressEvent(self, event):
        if event.button() == Qt.LeftButton:
            self.drag_pos = event.globalPos() - self.pos()
            event.accept()

    def mouseMoveEvent(self, event):
        if event.buttons() == Qt.LeftButton and self.drag_pos:
            self.move(event.globalPos() - self.drag_pos)
            event.accept()

    def resizeEvent(self, event):
        path = QPainterPath()
        path.addRoundedRect(QRectF(0, 0, self.width(), self.height()), 50, 50)
        region = QRegion(path.toFillPolygon().toPolygon())
        self.setMask(region)
        super().resizeEvent(event)

if __name__ == "__main__":
    app = QApplication(sys.argv)
    app.setWindowIcon(QIcon('./data_init/app_icon.ico'))
    window = IntegratedApp()
    window.show()
    sys.exit(app.exec_())