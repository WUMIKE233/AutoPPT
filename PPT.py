# 智能PPT生成器（v6.30）
#最近更新：UI优化完毕
import sys
import os
from datetime import datetime
import threading
import random
from PyQt5.QtWidgets import (QApplication, QMainWindow, QWidget, QVBoxLayout, QHBoxLayout, QLabel, QTextEdit,
                             QPushButton, QRadioButton, QFileDialog, QComboBox, QProgressBar, QFrame, QLineEdit, QMessageBox)
from PyQt5.QtGui import QPixmap, QBrush, QFont, QPalette, QIcon, QColor, QRadialGradient, QPainterPath, QRegion
from PyQt5.QtCore import Qt, QTimer, QThread, pyqtSignal, QRect, QRectF

from PIL import Image
import logging
import re
import requests
import time
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE

# 配置区域
DEEPSEEK_API_KEY = "sk-986dc7e18c9442aaac66207428d28600"
SAFE_CHARS = set("!@$%^&()_-+=[]{}',.`~")
TEMPLATE_MAP = {
    "教学课件类": "template_teaching.pptx",
    "政务报告类": "template_gov.pptx",
    "学术答辩类": "template_academic.pptx",
    "简约类": "template_simple.pptx"
}
MAX_RETRIES = 3
TIMEOUT = 80

# 日志配置
logging.basicConfig(
    filename=f'ppt_gen_{datetime.now().strftime("%Y%m%d")}.log',
    level=logging.DEBUG,
    format='%(asctime)s - %(levelname)s - %(message)s',
    encoding='utf-8'
)

# PPTGenerator 类（保持不变）
class PPTGenerator:
    def __init__(self, update_status=None):
        self.prs = None
        self.layout_map = {'cover': 0, 'section': 1, 'content': [2, 3], 'chart': 3, 'quote': 4, 'end': 4}
        self.available_content_layouts = [2, 3]
        self.cancelled = False
        self.update_status = update_status
        self.start_time = None
        self.current_stage = "就绪"
        self.selected_template = "template_simple.pptx"

    def set_template(self, template_name):
        self.selected_template = TEMPLATE_MAP.get(template_name, "template_simple.pptx")
        logging.info(f"用户选择模板: {template_name} -> {self.selected_template}")

    def safe_filename(self, filename):
        return ''.join(c for c in filename if c.isalnum() or c in SAFE_CHARS)

    def clean_markdown(self, text):
        text = re.sub(r'[#\*\>\-]+\s*', '', text)
        return text.strip()

    def call_deepseek_api(self, prompt):
        headers = {"Authorization": f"Bearer {DEEPSEEK_API_KEY}", "Content-Type": "application/json"}
        data = {"messages": [{"role": "user", "content": prompt}], "model": "deepseek-chat", "temperature": 0.5}
        for attempt in range(MAX_RETRIES):
            try:
                response = requests.post("https://api.deepseek.com/v1/chat/completions", headers=headers, json=data,
                                         timeout=TIMEOUT)
                response.raise_for_status()
                result = response.json()
                content = result['choices'][0]['message']['content']
                logging.debug(f"API返回内容长度: {len(content)} 字符, 内容: {content[:100]}...")
                return content[:1500] if len(content) > 1500 else content
            except requests.Timeout:
                logging.warning(f"API请求超时，第 {attempt + 1}/{MAX_RETRIES} 次尝试")
                if attempt == MAX_RETRIES - 1:
                    raise Exception("API请求超时，请检查网络连接")
                time.sleep(2 ** attempt)
            except Exception as e:
                logging.error(f"API调用失败: {str(e)}")
                raise Exception(f"API调用失败: {str(e)}")

    def generate_outline(self, topic):
        self.current_stage = "正在生成大纲"
        if self.update_status:
            self.update_status(self.current_stage)
        prompt = (
            f"生成关于【{topic}】的PPT大纲，格式为Markdown，包含4-6个章节，每个章节3-5个要点，"
            f"要求内容详细且丰富，提供具体的描述或示例，插入数据标记如{{data:market_size}}"
        )
        outline = self.call_deepseek_api(prompt)
        logging.info(f"生成的大纲长度: {len(outline)} 字符, 大纲: {outline}")
        return outline

    def expand_content(self, section_title, content):
        self.current_stage = f"正在扩展章节内容: {section_title}"
        if self.update_status:
            self.update_status(self.current_stage)
        content = re.sub(re.escape(section_title), '', content, flags=re.IGNORECASE).strip()
        prompt = (
            f"请为PPT章节【{section_title}】生成详细内容，基于以下现有内容：\n"
            f"```\n{content}\n```\n"
            f"严格按照以下Markdown格式返回：\n"
            f"```\n"
            f"## 正文\n"
            f"基于现有内容优化或补充，生成200-300字的详细正文，要求内容丰富，包含具体细节、示例或数据（如市场规模、趋势分析等），按段落分隔。\n"
            f"## 备注\n"
            f"生成与正文相关的拓展内容，必须包含以下两部分：\n"
            f"- 知识点：列出2-3个与正文相关的核心知识点（每点一句）。\n"
            f"- 提问：列出2-3个适合课堂提问的简单问题（每题一句）。\n"
            f"```\n"
            f"示例：\n"
            f"```\n"
            f"## 正文\n"
            f"机器学习是AI的核心技术，通过数据训练模型预测结果。例如，垃圾邮件过滤器利用分类算法，准确率达95%。\n"
            f"## 备注\n"
            f"- 知识点：\n"
            f"  - 机器学习基于统计学和概率论。\n"
            f"  - 深度学习是神经网络的一种。\n"
            f"- 提问：\n"
            f"  - 机器学习如何通过数据学习？\n"
            f"  - 深度学习有哪些应用？\n"
            f"```\n"
            f"确保备注部分不可为空，且与正文内容相关！"
        )
        try:
            expanded = self.call_deepseek_api(prompt)
            logging.info(f"扩展内容总长度: {len(expanded)} 字符, 原始返回: {expanded[:200]}...")
            main_content_match = re.search(r'## 正文\s*([\s\S]*?)(?=## 备注|$)', expanded, re.DOTALL)
            notes_content_match = re.search(r'## 备注\s*([\s\S]*)', expanded, re.DOTALL)
            main_content = main_content_match.group(1).strip() if main_content_match else content
            notes_content = notes_content_match.group(1).strip() if notes_content_match else None

            if not notes_content or "知识点" not in notes_content or "提问" not in notes_content:
                logging.warning(f"API返回的备注格式不完整: {notes_content}")
                notes_content = (
                    "知识点：\n"
                    "- 此处为默认知识点1。\n"
                    "- 此处为默认知识点2。\n"
                    "提问：\n"
                    "- 此处为默认问题1？\n"
                    "- 此处为默认问题2？"
                )

            main_content = re.sub(r'{data:\w+}', 'X', main_content)
            notes_content = self.clean_markdown(notes_content)
            logging.debug(f"正文长度: {len(main_content)} 字符, 备注长度: {len(notes_content)} 字符")
            return main_content, notes_content
        except Exception as e:
            logging.error(f"扩展内容失败: {str(e)}")
            default_main = re.sub(r'{data:\w+}', 'X', content)
            default_notes = (
                "知识点：\n"
                "- 此处为默认知识点1。\n"
                "- 此处为默认知识点2。\n"
                "提问：\n"
                "- 此处为默认问题1？\n"
                "- 此处为默认问题2？"
            )
            return default_main, default_notes

    def parse_input(self, input_text):
        self.current_stage = "正在解析输入"
        if self.update_status:
            self.update_status(self.current_stage)
        if not input_text.strip():
            raise ValueError("输入内容不能为空")
        if re.search(r'^#+\s|\* |- |\d+\.', input_text, re.MULTILINE):
            return self.parse_markdown(input_text)
        else:
            return self.parse_markdown(self.generate_outline(input_text))

    def parse_markdown(self, md_text):
        data = {"title": "", "sections": []}
        current_section = None
        current_paragraph = []

        for line in md_text.split('\n'):
            line = line.strip()
            if not line:
                if current_paragraph and current_section:
                    current_section["content"].append(self.clean_markdown('\n'.join(current_paragraph)))
                    current_paragraph = []
                continue

            if line.startswith('# '):
                data["title"] = self.clean_markdown(line)
                logging.debug(f"解析封面标题: {data['title']}")
            elif line.startswith('## '):
                if current_section:
                    if current_paragraph:
                        current_section["content"].append(self.clean_markdown('\n'.join(current_paragraph)))
                        current_paragraph = []
                    data["sections"].append(current_section)
                cleaned_title = self.clean_markdown(line)
                current_section = {"title": cleaned_title, "content": []}
                logging.debug(f"解析二级标题: {cleaned_title}")
            elif re.match(r'^[\-\*] ', line) and current_section:
                if current_paragraph:
                    current_section["content"].append(self.clean_markdown('\n'.join(current_paragraph)))
                    current_paragraph = []
                cleaned_content = self.clean_markdown(line)
                if cleaned_content and cleaned_content != current_section["title"]:
                    current_section["content"].append(cleaned_content)
            elif re.match(r'^>\s+', line) and current_section:
                if current_paragraph:
                    current_section["content"].append(self.clean_markdown('\n'.join(current_paragraph)))
                    current_paragraph = []
                current_section["content"].append(self.clean_markdown(line))
            elif current_section:
                current_paragraph.append(line)

        if current_section:
            if current_paragraph:
                current_section["content"].append(self.clean_markdown('\n'.join(current_paragraph)))
            data["sections"].append(current_section)

        if not data["title"] or not data["sections"]:
            raise ValueError("输入必须包含标题和至少一个章节")
        return data

    def build_ppt(self, data, output_path):
        if self.cancelled:
            logging.info("生成被用户取消")
            return None
        self.start_time = time.time()
        output_path = os.path.abspath(output_path)
        logging.info(f"构建PPT时传入的路径: {output_path}")
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        self._init_presentation()
        self._create_cover_slide(data)
        self._create_toc_slide(data)
        for idx, section in enumerate(data["sections"]):
            if self.cancelled:
                logging.info("生成被用户取消")
                return None
            layout_type = self._detect_layout_type(section)
            self._add_section_slides(layout_type, section, idx)
        self._add_end_slide()
        self._save_ppt(output_path)
        return output_path

    def _init_presentation(self):
        self.current_stage = "正在初始化PPT"
        if self.update_status:
            self.update_status(self.current_stage)
        try:
            self.prs = Presentation(self.selected_template)
            while self.prs.slides:
                self.prs.part.drop_rel(self.prs.slides._sldIdLst[-1].rId)
                del self.prs.slides._sldIdLst[-1]
        except Exception as e:
            logging.warning(f"加载模板 {self.selected_template} 失败: {str(e)}，使用默认设置")
            self.prs = Presentation()
            self.available_content_layouts = [2, 3]
            self.prs.slide_width = Inches(13.33)
            self.prs.slide_height = Inches(7.5)

    def _create_cover_slide(self, data):
        self.current_stage = "生成封面页"
        if self.update_status:
            self.update_status(self.current_stage)
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[self.layout_map['cover']])
        title = slide.shapes.title or slide.shapes.add_textbox(Inches(1), Inches(1), Inches(11), Inches(2))
        title.text = data["title"]
        p = title.text_frame.paragraphs[0]
        p.font.size = Pt(44)
        p.font.name = '微软雅黑'
        p.alignment = PP_ALIGN.CENTER
        subtitle = None
        if len(slide.placeholders) > 1:
            for ph in slide.placeholders:
                if ph.placeholder_format.idx != 0:
                    subtitle = ph
                    break
        if not subtitle:
            logging.info(f"封面页缺少副标题占位符，自动添加文本框")
            subtitle = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11), Inches(1))
        subtitle.text = "吴卓宪同学制作的AI智能体制作"
        p = subtitle.text_frame.paragraphs[0]
        p.font.size = Pt(24)
        p.font.name = '微软雅黑'
        p.font.color.rgb = RGBColor(0x42, 0x85, 0xF4)
        p.alignment = PP_ALIGN.CENTER
        logging.debug(f"封面页标题: {data['title']}")

    def _create_toc_slide(self, data):
        self.current_stage = "生成目录页"
        if self.update_status:
            self.update_status(self.current_stage)
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[2])
        title = slide.shapes.title or slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))
        title.text = "目录"
        p = title.text_frame.paragraphs[0]
        p.font.name = '微软雅黑'
        p.font.size = Pt(36)
        p.alignment = PP_ALIGN.CENTER

        content_box = None
        if len(slide.placeholders) > 1:
            for ph in slide.placeholders:
                if ph.placeholder_format.idx != 0:
                    content_box = ph
                    break
        if not content_box:
            logging.info(f"目录页缺少正文占位符，自动添加文本框")
            content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(4))

        tf = content_box.text_frame
        tf.clear()
        tf.word_wrap = True
        for idx, section in enumerate(data["sections"], 1):
            p = tf.add_paragraph()
            p.text = f"{idx}. {section['title']}"
            p.font.size = Pt(20)
            p.font.name = '微软雅黑'
            p.space_after = Pt(10)
            p.level = 0

    def _detect_layout_type(self, section):
        content = "\n".join(section.get("content", []))
        if "{chart:" in content:
            return 'chart'
        if re.search(r'^>\s+', content, re.MULTILINE):
            return 'quote'
        if len(section.get("content", [])) < 2:
            return 'section'
        return 'content'

    def _create_slide(self, layout_type):
        try:
            if layout_type == 'content':
                layout_idx = random.choice(self.available_content_layouts)
            else:
                layout_idx = self.layout_map[layout_type] if isinstance(self.layout_map[layout_type], int) else \
                self.layout_map[layout_type][0]
            slide = self.prs.slides.add_slide(self.prs.slide_layouts[layout_idx])
            logging.debug(f"创建幻灯片 {layout_type}，使用版式索引 {layout_idx}，占位符数: {len(slide.placeholders)}")
            return slide
        except (IndexError, KeyError) as e:
            logging.warning(f"布局 {layout_type} 不可用: {str(e)}，使用默认布局")
            return self.prs.slides.add_slide(self.prs.slide_layouts[1])

    def _add_section_slides(self, layout_type, section, idx):
        self.current_stage = f"生成章节: {idx + 1}. {section['title']}"
        if self.update_status:
            self.update_status(self.current_stage)
        content = "\n".join(section.get("content", [])) or f"这是关于 {section['title']} 的默认内容。"
        main_content, notes_content = self.expand_content(section['title'], content)
        logging.info(
            f"章节 {idx + 1}. {section['title']} 正文长度: {len(main_content)} 字符, 备注长度: {len(notes_content)} 字符")

        if layout_type in ['chart', 'quote']:
            slide = self._create_slide(layout_type)
            self._add_slide_content(slide, section, idx, layout_type, main_content, notes_content=notes_content)
        else:
            content_parts = self._split_content(main_content)
            for i, part in enumerate(content_parts, 1):
                if self.cancelled:
                    return
                slide = self._create_slide(layout_type)
                title_suffix = f" ({i}/{len(content_parts)})" if len(content_parts) > 1 else ""
                self._add_slide_content(slide, section, idx, layout_type, part, title_suffix,
                                        notes_content if i == 1 else "（续前页备注）")

    def _split_content(self, content):
        if len(content) <= 800:
            return [content]
        parts = []
        current_part = ""
        for line in content.split('\n'):
            if len(current_part) + len(line) + 1 > 800 and current_part:
                parts.append(current_part.strip())
                current_part = line
            else:
                current_part += f"\n{line}" if current_part else line
        if current_part:
            parts.append(current_part.strip())
        return parts

    def _add_slide_content(self, slide, section, idx, layout_type, content, title_suffix="", notes_content=None):
        title = slide.shapes.title or slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))
        title.text = f"{idx + 1}. {section['title']}{title_suffix}"
        p = title.text_frame.paragraphs[0]
        p.font.name = '微软雅黑'
        p.font.size = {'cover': Pt(44), 'section': Pt(36), 'content': Pt(32), 'chart': Pt(28), 'quote': Pt(24)}.get(layout_type, Pt(32))
        p.alignment = PP_ALIGN.CENTER if layout_type in ['section', 'quote'] else PP_ALIGN.LEFT

        if layout_type == 'chart':
            self._add_chart(slide, content)
        elif layout_type == 'quote':
            self._add_quote(slide, content)
        else:
            self._add_normal_content(slide, content, layout_type)

        if notes_content:
            try:
                notes_slide = slide.notes_slide
                notes_tf = notes_slide.notes_text_frame
                notes_tf.clear()
                notes_tf.text = notes_content
                for p in notes_tf.paragraphs:
                    p.font.size = Pt(12)
                    p.font.name = '微软雅黑'
                logging.debug(f"幻灯片 {slide.slide_id} 添加备注成功: {notes_content[:50]}...")
            except Exception as e:
                logging.error(f"添加备注失败: {str(e)}")
                notes_tf.text = "备注添加失败，请检查模板或内容"

    def _add_normal_content(self, slide, content, layout_type):
        content_box = None
        if len(slide.placeholders) > 1:
            for ph in slide.placeholders:
                if ph.placeholder_format.idx != 0:
                    content_box = ph
                    break
        if not content_box:
            logging.info(f"幻灯片 {slide.slide_id} 无正文占位符，自动添加更大文本框")
            content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5.5))

        tf = content_box.text_frame
        tf.clear()
        tf.word_wrap = True
        paragraphs = content.split('\n\n')
        for i, para in enumerate(paragraphs):
            if para.strip():
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = para.strip()
                p.font.size = Pt(18)
                p.font.name = '微软雅黑'
                p.space_after = Pt(8)
                if layout_type == 'content':
                    p.level = 0
        logging.debug(
            f"幻灯片 {slide.slide_id} 正文字符数: {len(content)}, 段落数: {len([p for p in paragraphs if p.strip()])}")

    def _add_quote(self, slide, content):
        content_box = None
        if len(slide.placeholders) > 1:
            for ph in slide.placeholders:
                if ph.placeholder_format.idx != 0:
                    content_box = ph
                    break
        if not content_box:
            content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5.5))

        tf = content_box.text_frame
        tf.clear()
        tf.text = content
        for p in tf.paragraphs:
            p.font.italic = True
            p.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
            p.alignment = PP_ALIGN.CENTER
        logging.debug(f"引用页正文: {content}")

    def _add_chart(self, slide, content):
        chart_type_match = re.search(r'{chart:(\w+)}', content)
        chart_type = chart_type_match.group(1) if chart_type_match else "default"
        chart_data = ChartData()
        chart_data.categories = ['2021', '2022', '2023']
        chart_data.add_series(chart_type.title(), (120, 180, 250) if chart_type == 'market_size' else (
        15, 25, 30) if chart_type == 'growth_rate' else (100, 200, 300))
        chart = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1), Inches(1.5), Inches(10), Inches(5), chart_data
        ).chart
        chart.has_title = True
        chart.chart_title.text_frame.text = f"{chart_type.replace('_', ' ').title()} 数据"
        logging.debug(f"图表页标题: {chart.chart_title.text_frame.text}")

    def _add_end_slide(self):
        self.current_stage = "生成结束页"
        if self.update_status:
            self.update_status(self.current_stage)
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[self.layout_map['end']])
        title = slide.shapes.title or slide.shapes.add_textbox(Inches(1), Inches(1), Inches(11), Inches(1))
        title.text = "感谢聆听"
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x2E, 0x86, 0xC1)

        content_box = None
        if len(slide.placeholders) > 1:
            for ph in slide.placeholders:
                if ph.placeholder_format.idx != 0:
                    content_box = ph
                    break
        if not content_box:
            content_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11), Inches(2))
        content_box.text_frame.text = "欢迎提问\ncontact@example.com"
        logging.debug(f"结束页标题: 感谢聆听")

    def _save_ppt(self, output_path):
        if not self.cancelled:
            self.current_stage = "正在保存PPT"
            if self.update_status:
                self.update_status(self.current_stage)
            output_path = os.path.abspath(output_path)
            logging.info(f"保存PPT时的最终路径: {output_path}")
            self.prs.save(output_path)
            if not os.path.exists(output_path):
                raise FileNotFoundError(f"文件未保存成功: {output_path}")
            logging.info(f"PPT保存成功至: {output_path}")

    def cancel(self):
        self.cancelled = True

class OutlineWorker(QThread):
    outline_generated = pyqtSignal(str)
    error = pyqtSignal(str)

    def __init__(self, generator, content):
        super().__init__()
        self.generator = generator
        self.content = content

    def run(self):
        try:
            if re.search(r'^#+\s|\* |- |\d+\.', self.content, re.MULTILINE):
                outline_text = self.content
            else:
                outline_text = self.generator.generate_outline(self.content)
            self.outline_generated.emit(outline_text)
        except Exception as e:
            self.error.emit(f"生成大纲失败: {str(e)}")

class PPTWorker(QThread):
    status_update = pyqtSignal(str)
    finished = pyqtSignal(str)
    error = pyqtSignal(str)

    def __init__(self, generator, outline_text, output_path):
        super().__init__()
        self.generator = generator
        self.outline_text = outline_text
        self.output_path = output_path
        self.cancelled = False

    def run(self):
        try:
            data = self.generator.parse_input(self.outline_text)
            path = self.generator.build_ppt(data, self.output_path)
            if path and not self.generator.cancelled:
                self.finished.emit(f"PPT已保存至:\n{path}")
            elif self.generator.cancelled:
                self.finished.emit("PPT生成已被取消")
        except Exception as e:
            self.error.emit(f"发生错误: {str(e)}\n建议: 检查输入内容或文件夹权限。")

    def cancel(self):
        self.generator.cancel()
        self.cancelled = True

class PPTApp(QMainWindow):
    def __init__(self):
        super().__init__()
        self.running = False
        self.selected_folder = "D:\\outputPPT"
        self.generator = PPTGenerator(self.update_status)
        self.worker = None
        self.outline_worker = None
        self.start_time = None
        self.outline_text = ""
        self.logo_path = "./data_init/logo.png"
        self.setWindowFlags(Qt.FramelessWindowHint)
        self.setAttribute(Qt.WA_TranslucentBackground, True)
        self.init_ui()

    def init_ui(self):
        self.setWindowTitle("智能PPT生成系统 v6.50")
        self.setGeometry(100, 100, 1200, 900)

        # 中央部件设置渐变背景和圆角
        central_widget = QWidget(self)
        central_widget.setStyleSheet("""
            background-color: qradialgradient(cx: 0%, cy: 100%, radius: 100%, stop: 0 #00CED1, stop: 1 #87CEEB);
            border-radius: 50px;
            border: 1px solid #AAAAAA;
        """)
        self.setCentralWidget(central_widget)

        # 主布局
        main_layout = QVBoxLayout(central_widget)
        main_layout.setContentsMargins(20, 20, 20, 20)
        main_layout.setSpacing(15)

        # 标题栏
        title_bar = QWidget(central_widget)
        title_bar.setFixedHeight(140)
        title_bar.setStyleSheet("background: transparent; border: none;")
        title_bar_layout = QHBoxLayout(title_bar)
        title_bar_layout.setContentsMargins(0, 0, 0, 0)

        # Logo
        self.logo_label = QLabel(title_bar)
        self.logo_label.setFixedSize(197, 139)
        self.update_logo()

        # 标题
        self.title_label = QLabel("智能PPT生成系统 v6.50", title_bar)
        self.title_label.setFont(QFont("微软雅黑", 24, QFont.Bold))
        self.title_label.setStyleSheet("color: #333333; background: transparent;")
        self.title_label.setAlignment(Qt.AlignCenter)

        # 关闭按钮
        close_btn = QPushButton("×", title_bar)
        close_btn.setFont(QFont("微软雅黑", 14, QFont.Bold))
        close_btn.setFixedSize(30, 30)
        close_btn.setStyleSheet("""
            QPushButton {
                background: transparent;
                color: #666666;
                border: none;
            }
            QPushButton:hover {
                color: #FF0000;
            }
        """)
        close_btn.clicked.connect(self.close)

        title_bar_layout.addWidget(self.logo_label)
        title_bar_layout.addWidget(self.title_label)
        title_bar_layout.addStretch()
        title_bar_layout.addWidget(close_btn)
        main_layout.addWidget(title_bar)

        # 输入方式选择，设置框架背景为透明
        input_frame = QFrame(central_widget)
        input_frame.setStyleSheet("background: transparent;")
        input_layout = QHBoxLayout(input_frame)
        self.input_text_radio = QRadioButton("文本输入", input_frame)
        self.input_file_radio = QRadioButton("文件上传", input_frame)
        self.input_text_radio.setChecked(True)
        self.input_text_radio.toggled.connect(self.toggle_input)
        for radio in [self.input_text_radio, self.input_file_radio]:
            radio.setFont(QFont("微软雅黑", 12))
        input_layout.addWidget(self.input_text_radio)
        input_layout.addWidget(self.input_file_radio)
        main_layout.addWidget(input_frame)

        # 文本输入区
        self.text_input = QTextEdit(central_widget)
        self.text_input.setFont(QFont("微软雅黑", 12))
        main_layout.addWidget(self.text_input)

        # 文件上传区，设置框架背景为透明
        self.file_frame = QFrame(central_widget)
        self.file_frame.setStyleSheet("background: transparent;")
        file_layout = QHBoxLayout(self.file_frame)
        file_upload_label = QLabel("文件上传:", self.file_frame)
        file_upload_label.setFont(QFont("微软雅黑", 12))
        self.file_path = QLineEdit(self.file_frame)
        self.file_path.setReadOnly(True)
        self.file_path.setFixedHeight(60)
        file_btn = QPushButton("选择文件", self.file_frame)
        file_btn.setFixedHeight(60)
        file_btn.setFont(QFont("微软雅黑", 12))
        file_btn.clicked.connect(self.load_file)
        file_layout.addWidget(file_upload_label)
        file_layout.addWidget(self.file_path)
        file_layout.addWidget(file_btn)
        self.file_frame.hide()
        main_layout.addWidget(self.file_frame)

        # 模板选择，设置框架背景为透明
        template_frame = QFrame(central_widget)
        template_frame.setStyleSheet("background: transparent;")
        template_layout = QHBoxLayout(template_frame)
        template_label = QLabel("选择模板:", template_frame)  # 修复父部件为 template_frame
        template_label.setFont(QFont("微软雅黑", 12))
        self.template_combo = QComboBox(template_frame)  # 修复父部件为 template_frame
        self.template_combo.addItems(TEMPLATE_MAP.keys())
        self.template_combo.setCurrentText("简约类")
        self.template_combo.setFixedHeight(60)
        self.template_combo.setFont(QFont("微软雅黑", 12))
        self.template_combo.currentTextChanged.connect(self.update_template)
        # 设置模板下拉表的样式
        self.template_combo.setStyleSheet("""
            QComboBox {
                background-color: rgba(255, 255, 255, 0.9);
                color: #333333;
                border: 1px solid #AAAAAA;
                border-radius: 10px;
                padding: 5px;
            }
            QComboBox::drop-down {
                border: none;
                width: 20px;
            }
            QComboBox::down-arrow {
                image: url(./data_init/down_arrow.png);  # 可选：添加下拉箭头图标
                width: 10px;
                height: 10px;
            }
            QComboBox QAbstractItemView {
                background-color: rgba(255, 255, 255, 0.95);
                color: #333333;
                selection-background-color: rgba(53, 142, 255, 0.8);
                selection-color: white;
                border: 1px solid #AAAAAA;
                border-radius: 1px;
            }
        """)
        template_layout.addWidget(template_label)
        template_layout.addWidget(self.template_combo)
        main_layout.addWidget(template_frame)

        # 输出设置，设置框架背景为透明
        output_frame = QFrame(central_widget)
        output_frame.setStyleSheet("background: transparent;")
        output_layout = QHBoxLayout(output_frame)
        folder_label = QLabel("输出文件夹:", output_frame)  # 修复父部件为 output_frame
        folder_label.setFont(QFont("微软雅黑", 12))
        self.folder_input = QLineEdit(self.selected_folder, output_frame)
        self.folder_input.setFixedHeight(60)
        folder_btn = QPushButton("选择文件夹", output_frame)  # 修复父部件为 output_frame
        folder_btn.setFixedHeight(60)
        folder_btn.setFont(QFont("微软雅黑", 12))
        folder_btn.clicked.connect(self.select_folder)
        output_label = QLabel("输出文件:", output_frame)  # 修复父部件为 output_frame
        output_label.setFont(QFont("微软雅黑", 12))
        self.output_path = QLineEdit(self.get_default_output_path(), output_frame)
        self.output_path.setReadOnly(True)
        self.output_path.setFixedHeight(60)
        output_layout.addWidget(folder_label)
        output_layout.addWidget(self.folder_input)
        output_layout.addWidget(folder_btn)
        output_layout.addWidget(output_label)
        output_layout.addWidget(self.output_path)
        main_layout.addWidget(output_frame)

        # 大纲显示
        self.outline_display = QTextEdit(central_widget)
        self.outline_display.setFont(QFont("微软雅黑", 12))
        main_layout.addWidget(self.outline_display)

        # 控制按钮区，设置框架背景为透明
        control_frame = QFrame(central_widget)
        control_frame.setStyleSheet("background: transparent;")
        control_layout = QHBoxLayout(control_frame)
        self.progress = QProgressBar(control_frame)  # 修复父部件为 control_frame
        self.progress.setRange(0, 0)
        self.progress.hide()
        self.preview_btn = QPushButton("生成大纲", control_frame)  # 修复父部件为 control_frame
        self.preview_btn.setFixedHeight(60)
        self.preview_btn.setFont(QFont("微软雅黑", 12, QFont.Bold))
        self.preview_btn.clicked.connect(self.preview_outline)
        self.generate_btn = QPushButton("确认生成PPT", control_frame)  # 修复父部件为 control_frame
        self.generate_btn.setFixedHeight(60)
        self.generate_btn.setFont(QFont("微软雅黑", 12, QFont.Bold))
        self.generate_btn.clicked.connect(self.start_generate)
        self.generate_btn.setEnabled(False)
        self.cancel_btn = QPushButton("取消", control_frame)  # 修复父部件为 control_frame
        self.cancel_btn.setFixedHeight(60)
        self.cancel_btn.setFont(QFont("微软雅黑", 12, QFont.Bold))
        self.cancel_btn.clicked.connect(self.cancel_generate)
        self.cancel_btn.setEnabled(False)
        self.status_label = QLabel("就绪", control_frame)  # 修复父部件为 control_frame
        self.status_label.setStyleSheet("color: white; background: transparent;")
        self.status_label.setFont(QFont("微软雅黑", 12))
        control_layout.addWidget(self.progress)
        control_layout.addWidget(self.preview_btn)
        control_layout.addWidget(self.generate_btn)
        control_layout.addWidget(self.cancel_btn)
        control_layout.addWidget(self.status_label)
        main_layout.addWidget(control_frame)

        # 统一按钮样式，与 MainApp 一致
        for btn in [file_btn, folder_btn, self.preview_btn, self.generate_btn, self.cancel_btn]:
            btn.setStyleSheet("""
                QPushButton {
                    background: rgba(53, 142, 255, 0.9); 
                    color: white; 
                    border: none; 
                    border-radius: 10px; 
                    padding: 10px 20px;
                }
                QPushButton:disabled {background: rgba(150, 150, 150, 0.6);}
                QPushButton:hover {background: rgba(53, 142, 255, 1);}
            """)

        # 添加动态标题颜色效果
        self.update_title_color()

    def resizeEvent(self, event):
        super().resizeEvent(event)
        path = QPainterPath()
        path.addRoundedRect(QRectF(0, 0, self.width(), self.height()), 50, 50)
        region = QRegion(path.toFillPolygon().toPolygon())
        self.setMask(region)

    def update_logo(self):
        if os.path.exists(self.logo_path):
            pixmap = QPixmap(self.logo_path)
            pixmap = pixmap.scaled(self.logo_label.size(), Qt.KeepAspectRatio, Qt.SmoothTransformation)
            self.logo_label.setPixmap(pixmap)
            self.logo_label.setStyleSheet("background: transparent;")
        else:
            logging.warning(f"Logo文件不存在: {self.logo_path}")

    def update_title_color(self):
        r = random.randint(0, 255)
        g = random.randint(0, 255)
        b = random.randint(0, 255)
        self.title_label.setStyleSheet(f"color: rgb({r}, {g}, {b}); background: transparent;")
        QTimer.singleShot(500, self.update_title_color)

    def mousePressEvent(self, event):
        if event.button() == Qt.LeftButton:
            self.drag_pos = event.globalPos() - self.pos()
            event.accept()

    def mouseMoveEvent(self, event):
        if event.buttons() == Qt.LeftButton and self.drag_pos:
            self.move(event.globalPos() - self.drag_pos)
            event.accept()

    def toggle_input(self):
        if self.input_text_radio.isChecked():
            self.text_input.show()
            self.file_frame.hide()
        else:
            self.text_input.hide()
            self.file_frame.show()

    def load_file(self):
        path, _ = QFileDialog.getOpenFileName(self, "选择文件", "", "Markdown文件 (*.md)")
        if path:
            self.file_path.setText(path)
            with open(path, 'r', encoding='utf-8') as f:
                self.text_input.setPlainText(f.read())

    def load_file_content(self):
        try:
            with open(self.file_path.text(), 'r', encoding='utf-8') as f:
                return f.read()
        except Exception as e:
            logging.error(f"加载文件失败: {str(e)}")
            self.status_label.setText(f"错误: 无法加载文件")
            return ""

    def on_outline_generated(self, outline_text):
        self.outline_display.setPlainText(outline_text)
        self.generate_btn.setEnabled(True)
        self.running = False
        self.progress.hide()
        self.status_label.setText("就绪")
        self.preview_btn.setEnabled(True)

    def on_outline_error(self, error_message):
        logging.error(error_message)
        self.status_label.setText(error_message)
        self.running = False
        self.progress.hide()
        self.preview_btn.setEnabled(True)
        self.generate_btn.setEnabled(False)
        QMessageBox.critical(self, "错误", error_message)

    def select_folder(self):
        folder = QFileDialog.getExistingDirectory(self, "选择文件夹", self.selected_folder)
        if folder:
            self.selected_folder = folder
            self.folder_input.setText(folder)
            self.output_path.setText(self.get_default_output_path())

    def get_default_output_path(self):
        folder = os.path.abspath(self.selected_folder)
        filename = datetime.now().strftime("%Y%m%d%H%M") + ".pptx"
        output_path = os.path.join(folder, filename)
        os.makedirs(folder, exist_ok=True)
        return output_path

    def update_template(self, template_name):
        self.generator.set_template(template_name)

    def preview_outline(self):
        if not self.running:
            self.running = True
            self.progress.show()
            self.status_label.setText("生成大纲中...")
            self.preview_btn.setEnabled(False)
            self.generate_btn.setEnabled(False)
            content = self.text_input.toPlainText() if self.input_text_radio.isChecked() else self.load_file_content()
            self.outline_worker = OutlineWorker(self.generator, content)
            self.outline_worker.outline_generated.connect(self.on_outline_generated)
            self.outline_worker.error.connect(self.on_outline_error)
            self.outline_worker.start()

    def start_generate(self):
        if not self.running:
            self.running = True
            self.start_time = datetime.now()
            self.progress.show()
            self.status_label.setText("生成PPT中...")
            self.generate_btn.setEnabled(False)
            self.cancel_btn.setEnabled(True)
            self.output_path.setText(self.get_default_output_path())
            self.outline_text = self.outline_display.toPlainText()
            self.worker = PPTWorker(self.generator, self.outline_text, self.output_path.text())
            self.worker.status_update.connect(self.update_status)
            self.worker.finished.connect(self.on_generate_finished)
            self.worker.error.connect(self.on_generate_error)
            self.worker.start()
            self.update_ui()

    def cancel_generate(self):
        if self.running and self.worker:
            self.worker.cancel()
            self.status_label.setText("取消中...")

    def on_generate_finished(self, message):
        self.running = False
        self.progress.hide()
        self.status_label.setText("就绪")
        self.preview_btn.setEnabled(True)
        self.generate_btn.setEnabled(True)
        self.cancel_btn.setEnabled(False)
        QMessageBox.information(self, "成功", message)

    def on_generate_error(self, message):
        self.running = False
        self.progress.hide()
        self.status_label.setText("就绪")
        self.preview_btn.setEnabled(True)
        self.generate_btn.setEnabled(True)
        self.cancel_btn.setEnabled(False)
        QMessageBox.critical(self, "错误", message)

    def update_status(self, stage):
        self.generator.current_stage = stage

    def update_ui(self):
        if self.running and self.start_time:
            elapsed = (datetime.now() - self.start_time).total_seconds()
            self.status_label.setText(f"{self.generator.current_stage} (耗时: {elapsed:.2f}秒)")
            QTimer.singleShot(100, self.update_ui)
        elif not self.running:
            self.status_label.setText("就绪")

if __name__ == "__main__":
    app = QApplication(sys.argv)

    # 设置应用程序图标
    # 将 'app_icon.ico' 替换为你实际的图标文件路径
    app.setWindowIcon(QIcon('app_icon.ico'))  # 例如：'icons/my_icon.ico' 或绝对路径如 'C:/path/to/app_icon.ico'

    window = PPTApp()
    window.show()
    sys.exit(app.exec_())